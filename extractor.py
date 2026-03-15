"""Vision LLM extractor – sends page images to GPT-4o, Claude, or Kimi and returns structured JSON.

Supports two modes:
  1. Exam paper extraction (original) – extracts questions from past papers
  2. Notebook generation – generates an intuitive study guide from lecture notes/slides

Supports chunked parallel processing for large lecture decks.
"""

from __future__ import annotations

import base64
import concurrent.futures
import io
import json
import math
import os
import re
from typing import Any

from PIL import Image

# ---------------------------------------------------------------------------
# System prompt – instructs the model exactly how to produce the JSON output
# ---------------------------------------------------------------------------

SYSTEM_PROMPT = r"""You are an expert OCR and exam-paper analysis system.
You will receive one or more images of an exam / past paper. Your job is to
extract **every** question and return a single JSON object that conforms
exactly to the schema below.

### Rules
1. **Question types**
   - `"multiple-choice"` – has a fixed set of answer options (A/B/C/D …).
   - `"open-ended"` – free-response, essay, short-answer, or calculation.

2. **Math / equations**
   - If a question contains a mathematical formula, set `"equation"` to a
     LaTeX string (use standard LaTeX notation, e.g. `E_d = \frac{\%\Delta Q}{\%\Delta P}`).
   - If no formula is present, set `"equation"` to `null`.

3. **Diagrams / images**
   - If a question refers to or contains a diagram, graph, figure, or image,
     add a field `"images"` which is a list of short descriptive filenames
     like `"q3_diagram.png"`.  Also add a field `"imageDescriptions"` with
     a detailed textual description of each diagram so the content is
     preserved even without the image file.
   - If no diagram is present, set `"images"` to `null`.

4. **Correct answers**
   - For multiple-choice: set `"correctAnswerId"` to the letter of the
     correct option **only if you are confident** (e.g. answer key is
     visible or the answer is unambiguous). Otherwise set it to `null`.
   - For open-ended: provide a `"modelAnswer"` (a concise ideal answer)
     and `"keyTerms"` (list of important terms/concepts expected).

5. **IDs**
   - `"id"` of the paper: generate a snake_case slug from the title.
   - `"id"` of each question: use the pattern `"q1"`, `"q2"`, etc.

6. **Preserve original text** – transcribe question text as faithfully as
   possible, including any bold/italic markers if evident.

7. **Sub-questions** – if a question has parts (a), (b), (c), treat each
   part as a separate question with `"number"` reflecting hierarchy,
   e.g. `3` for question 3, and include "Part (a):" prefix in the text.

### Output schema (TypeScript-style for clarity)

```
{
  "id": string,
  "title": string,
  "description": string | null,
  "questions": Array<MCQ | OpenEnded>
}

MCQ = {
  "id": string,
  "number": number,
  "type": "multiple-choice",
  "text": string,
  "equation": string | null,
  "images": string[] | null,
  "imageDescriptions": string[] | null,
  "options": Array<{ "id": string, "text": string }>,
  "correctAnswerId": string | null
}

OpenEnded = {
  "id": string,
  "number": number,
  "type": "open-ended",
  "text": string,
  "equation": string | null,
  "images": string[] | null,
  "imageDescriptions": string[] | null,
  "modelAnswer": string | null,
  "keyTerms": string[] | null
}
```

Return **only** the JSON object. No markdown fences, no commentary.
"""

# ---------------------------------------------------------------------------
# Notebook generation prompt – produces a Socratic study guide from slides/notes
# ---------------------------------------------------------------------------

NOTEBOOK_SYSTEM_PROMPT = r"""You are an expert educational content designer and Socratic tutor.
You will receive one or more images of university lecture slides or notes. Your job is to
transform this raw lecture content into an **intuitive, comprehensive study guide** that
any student could follow — even without the original slides.

### Your Teaching Philosophy
- Explain concepts as if talking to a smart friend who's new to the topic.
- Use analogies, real-world examples, and "why does this matter?" framing.
- Make abstract ideas concrete. If there's a formula, explain the intuition BEFORE the math.
- Use a warm, encouraging tone — never condescending.
- Build understanding progressively: start with the big picture ("why does this exist?"),
  then dive into mechanics ("how does it work?"), then edge cases ("when does it break?").
- Include "Key Insight" callouts in subsections — one sentence that captures the core takeaway
  a student should remember even if they forget everything else.
- Where applicable, include a brief worked example or mini-problem with solution to ground
  the theory in practice. Students learn best by seeing concepts applied, not just described.

### Rules

1. **Structure**: Organize content into logical sections and subsections.
   Each section should cover one major concept from the lecture.

2. **Section titles**: Use the pattern "Topic Name (Intuitive Subtitle)" where the
   subtitle gives a hint about what the concept really means.
   Examples: "Price Elasticity (How Sensitive Are Buyers?)", "Derivatives (The Rate Detective)"

3. **Tags**: Each section MUST include `"tags"` — an array of lowercase keywords that
   describe the topics covered. These will be used to match against past paper questions.
   Be thorough: include the main concept, related terms, and common exam keywords.
   Example: ["elasticity", "price elasticity", "demand", "percentage change", "elastic", "inelastic"]

4. **Content**: The opening paragraph of each section should frame the concept intuitively.
   Subsections go deeper with explanations, worked examples, and key takeaways as bullets.

5. **Icons**: Choose a single emoji that represents each section's topic.

6. **Intro**: Write a welcoming introduction that frames the entire lecture and motivates
   the student. Include 2-4 "highlights" (key topics they'll learn).

7. **Chat responses**: Generate 5-8 keyword-triggered responses for an AI tutor chatbot.
   Each has `keywords` (pipe-separated, lowercase) and a `response` (helpful, Socratic).

8. **Metadata**: Generate a sensible `id` (slug), `title`, `course` name, `icon` (emoji),
   and `color` (hex code that matches the subject — e.g. green for economics, blue for math).

### COMPLETENESS RULES (CRITICAL — do not lose information)

9. **NEVER omit information** from the slides. If a slide contains a formula, definition,
   example, list, key term, or diagram description, it MUST appear in the output. When in
   doubt, include it. Being too detailed is always better than losing information.

10. **Formulas**: Include the formula itself, an intuitive plain-English explanation of what
    each variable means, and a worked example if one appears on the slides. Never skip a formula.

11. **Lists**: If a slide contains a list of items (e.g. "types of data structures",
    "properties of X"), include ALL items. Never summarize with "and more", "etc.", or
    "additional types exist". Write out every single item.

12. **Diagrams and figures**: If a slide references or contains a diagram, graph, or figure,
    describe what it shows, what the axes/labels represent, and what insight it conveys.

13. **Depth**: Each section should have at MINIMUM 3 subsections. Each subsection should have
    at minimum 2-3 bullet points. If the source material is dense, create more subsections
    rather than cramming everything into one.

14. **Definitions**: Every technical term introduced in the slides must have a clear definition
    in the notes. Do not assume the student already knows domain-specific vocabulary.

### Output Schema

```
{
  "id": string,                    // e.g. "nb_microecon"
  "title": string,                 // e.g. "Microeconomics: How Markets Think"
  "course": string,                // e.g. "Economics 101 — Microeconomics"
  "icon": string,                  // single emoji
  "color": string,                 // hex color, e.g. "#2ECC71"
  "intro": {
    "text": string,                // welcoming intro paragraph
    "highlights": [
      { "label": string, "desc": string }   // 2-4 key topics
    ]
  },
  "sections": [
    {
      "icon": string,              // single emoji
      "title": string,             // "Topic (Intuitive Subtitle)"
      "tags": string[],            // lowercase topic tags
      "content": string,           // opening paragraph
      "subsections": [             // optional deeper dives
        {
          "title": string,
          "content": string,
          "bullets": string[] | null
        }
      ]
    }
  ],
  "chatResponses": [
    {
      "keywords": string,          // "elasticity|elastic|inelastic"
      "response": string           // helpful Socratic response
    }
  ]
}
```

Return **only** the JSON object. No markdown fences, no commentary.
"""

# ---------------------------------------------------------------------------
# Question-matching prompt – matches past paper questions to notebook sections
# ---------------------------------------------------------------------------

MATCHING_SYSTEM_PROMPT = r"""You are an educational content matching system with deep domain knowledge.

You will receive:
1. Notebook sections with their titles, tags, and content summaries
2. A list of past paper questions from one or more papers

Your job is to match questions to notebook sections based on **semantic topic overlap** —
not just exact keyword matches.

### Semantic Matching Rules
- "supply and demand" matches "market equilibrium", "price mechanism", "demand curve"
- "elasticity" matches "price sensitivity", "percentage change in quantity"
- "derivative" matches "rate of change", "marginal", "differentiation"
- "probability" matches "likelihood", "chance", "random events"
- Tags may use different vocabulary — match by MEANING, not just exact words.
- A question about "consumer surplus under price controls" matches a section on "price ceilings"
  even if the tags don't share exact words.

### Selection Rules
- Only include questions that are genuinely relevant to the notebook's topics.
- A question is relevant if its content, concepts, or required knowledge overlaps with
  any section's topic area.
- Return IDs in order of relevance (most relevant first).
- Aim for 4-10 matched questions per notebook.
- If fewer than 3 questions match well, include the closest matches anyway.
- Prefer a mix of multiple-choice and open-ended questions when possible.

### Output format
Return a JSON array of question ID strings:
["paper_id_q1", "paper_id_q3", "paper_id_e2"]

Return **only** the JSON array. No markdown fences, no commentary.
"""


# ---------------------------------------------------------------------------
# Merge prompt – combines partial chunk outputs into one cohesive notebook
# ---------------------------------------------------------------------------

MERGE_SYSTEM_PROMPT = r"""You are an expert educational content editor. You will receive multiple JSON
fragments — each is a partial study guide generated from a batch of lecture slides. Your job is to
**merge** them into ONE cohesive, comprehensive study guide.

### Merge Rules

1. **Combine sections** — concatenate all sections into a single list. If two chunks covered the
   same topic (similar titles/tags), merge them into one richer section rather than duplicating.

2. **Deduplicate** — remove duplicate sections, subsections, or bullet points. Prefer the more
   detailed version when two are similar.

3. **Narrative flow** — reorder sections so they follow a logical teaching progression.
   Earlier/foundational concepts should come before advanced ones.

4. **Unified intro** — write ONE welcoming intro that frames the entire lecture holistically.
   Combine the highlights from all chunks, picking the 3-4 most important.

5. **Unified metadata** — pick the best `id`, `title`, `course`, `icon`, and `color` from
   the chunks. The title should capture the full lecture scope.

6. **Tags** — merge and deduplicate tags across all chunks. Be thorough.

7. **Chat responses** — merge all chatResponses, removing duplicates. Aim for 6-10 total.

8. **Depth** — each section should have at least 2-3 subsections. If a chunk's section is
   shallow, keep it but note that the content should remain as-is (don't invent content
   that wasn't in the slides).

### Output format
Return a SINGLE JSON object following the exact same Notebook schema as the inputs.
Return **only** the JSON object. No markdown fences, no commentary.

**CRITICAL**: You MUST output complete, valid JSON. Do NOT let the output get truncated.
If the merged content is very large, prioritize keeping all sections but write more concisely
within each subsection rather than risk an incomplete JSON response.
"""


# ---------------------------------------------------------------------------
# Hybrid extraction – text + embedded images (PPTX and digital PDFs)
# ---------------------------------------------------------------------------

HYBRID_NOTEBOOK_PROMPT = r"""You are an expert educational content designer and Socratic tutor.
You will receive the **extracted text content** from university lecture slides, and possibly some
embedded images (charts, diagrams, figures) that appeared in those slides.

Transform this content into an intuitive, comprehensive study guide following the EXACT same
rules and output schema as below.

### Your Teaching Philosophy
- Explain concepts as if talking to a smart friend who's new to the topic.
- Use analogies, real-world examples, and "why does this matter?" framing.
- Make abstract ideas concrete. If there's a formula, explain the intuition BEFORE the math.
- Use a warm, encouraging tone — never condescending.
- Build understanding progressively: start with the big picture ("why does this exist?"),
  then dive into mechanics ("how does it work?"), then edge cases ("when does it break?").
- Include "Key Insight" callouts in subsections — one sentence that captures the core takeaway
  a student should remember even if they forget everything else.
- Where applicable, include a brief worked example or mini-problem with solution to ground
  the theory in practice. Students learn best by seeing concepts applied, not just described.

### Rules
1. **Structure**: Organize content into logical sections and subsections.
2. **Section titles**: "Topic Name (Intuitive Subtitle)" pattern.
3. **Tags**: Each section MUST include "tags" — array of lowercase keywords for question matching.
4. **Content**: Opening paragraph frames the concept intuitively. Subsections go deeper.
5. **Icons**: Single emoji per section.
6. **Intro**: Welcoming introduction with 2-4 highlights.
7. **Chat responses**: 5-8 keyword-triggered responses for AI tutor chatbot.
8. **Metadata**: id (slug), title, course, icon (emoji), color (hex).

### COMPLETENESS RULES (CRITICAL)
9. NEVER omit information. Every formula, definition, example, list item must appear.
10. Formulas: include formula, plain-English explanation, worked example.
11. Lists: include ALL items. Never use "etc." or "and more".
12. **DIAGRAM PLACEMENT**: If embedded images are provided, place the marker `[DIAGRAM:N]` (where N is 1, 2, 3, … matching the order of images attached) at the exact spot in the content or subsection content where the diagram is most relevant. Also write a brief sentence before or after explaining what the diagram shows.
13. Each section: minimum 3 subsections, each with 2-3 bullet points.
14. Every technical term must have a clear definition.

### Output Schema
```
{
  "id": string, "title": string, "course": string, "icon": string, "color": string,
  "intro": { "text": string, "highlights": [{ "label": string, "desc": string }] },
  "sections": [
    {
      "icon": string, "title": string, "tags": string[],
      "content": string,
      "subsections": [{ "title": string, "content": string, "bullets": string[] | null }]
    }
  ],
  "chatResponses": [{ "keywords": string, "response": string }]
}
```

Return **only** the JSON object. No markdown fences, no commentary.
"""

MIN_TEXT_PER_PAGE = 50      # chars — below this, page is likely scanned/image-only
MIN_IMAGE_DIM = 150         # pixels — basic size gate
MAX_IMAGES_PER_CHUNK = 3    # cap vision tokens per API call
HYBRID_CHUNK_SIZE = 5       # smaller chunks for hybrid to stay under token limits
HYBRID_IMAGE_MAX_PX = 400   # resize embedded images to this max dimension
MIN_SLIDE_TEXT_FOR_IMAGES = 80  # chars — slides with less text are title/poll slides


def _image_phash(img: Image.Image) -> str:
    """Compute a simple perceptual hash (16x16 grayscale average hash)."""
    import numpy as np
    thumb = img.resize((16, 16), Image.LANCZOS).convert("L")
    pixels = np.array(thumb).flatten()
    avg = pixels.mean()
    return "".join("1" if p > avg else "0" for p in pixels)


def _basic_image_filter(img: Image.Image) -> bool:
    """Lightweight per-image filter: only reject tiny or extremely elongated images."""
    w, h = img.size
    if w < MIN_IMAGE_DIM or h < MIN_IMAGE_DIM:
        return False
    if w * h < 30_000:
        return False
    aspect = max(w, h) / max(min(w, h), 1)
    if aspect > 5.0:
        return False
    return True


def _filter_slides_images(slides_data: list[dict]) -> list[dict]:
    """Multi-pass image filtering using slide context and cross-slide dedup.

    Strategy:
      1. Basic per-image size/shape filter.
      2. Context filter: only keep images from slides with meaningful text.
         Slides with little text (title slides, poll slides, "any questions?")
         almost never contain educational diagrams.
      3. Cross-slide deduplication: images appearing on 2+ slides are decorative
         (logos, template icons, QR codes).
    """
    total_raw = sum(len(s["images"]) for s in slides_data)

    # Pass 1: basic size/shape
    for slide in slides_data:
        slide["images"] = [img for img in slide["images"] if _basic_image_filter(img)]
    after_basic = sum(len(s["images"]) for s in slides_data)

    # Pass 2: context — drop images from slides with little educational text
    for slide in slides_data:
        if slide["images"] and len(slide["text"]) < MIN_SLIDE_TEXT_FOR_IMAGES:
            slide["images"] = []
    after_context = sum(len(s["images"]) for s in slides_data)

    # Pass 3: cross-slide deduplication
    hash_to_slides: dict[str, set[int]] = {}
    for s_idx, slide in enumerate(slides_data):
        for img in slide["images"]:
            h = _image_phash(img)
            hash_to_slides.setdefault(h, set()).add(s_idx)

    repeated = {h for h, sids in hash_to_slides.items() if len(sids) >= 2}

    if repeated:
        for slide in slides_data:
            slide["images"] = [
                img for img in slide["images"]
                if _image_phash(img) not in repeated
            ]
    after_dedup = sum(len(s["images"]) for s in slides_data)

    print(
        f"  [filter] images: {total_raw} raw → {after_basic} size "
        f"→ {after_context} context → {after_dedup} dedup"
    )
    return slides_data


def extract_content_from_pptx(path: str) -> list[dict]:
    """Extract text and embedded images from a PPTX file.

    Returns a list of dicts per slide:
        [{ "slide_number": int, "text": str, "images": [PIL.Image] }]
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        texts = []
        images = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        texts.append(" | ".join(row_texts))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_blob = shape.image.blob
                    img = Image.open(io.BytesIO(img_blob)).convert("RGB")
                    images.append(img)
                except Exception:
                    pass

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Speaker notes: {notes}]")

        slides_data.append({
            "slide_number": idx + 1,
            "text": "\n".join(texts),
            "images": images,
        })

    slides_data = _filter_slides_images(slides_data)
    return slides_data


def extract_content_from_pdf(path: str) -> list[dict] | None:
    """Try to extract text + embedded images from a PDF.

    Returns slide data like extract_content_from_pptx, or None if the PDF
    appears to be scanned (insufficient text layer).
    """
    import fitz
    import pdfplumber

    slides_data = []
    total_text_len = 0

    pdf_plumber = pdfplumber.open(path)
    pdf_fitz = fitz.open(path)

    for page_idx in range(len(pdf_plumber.pages)):
        page_text = pdf_plumber.pages[page_idx].extract_text() or ""
        total_text_len += len(page_text)

        images = []
        if page_idx < len(pdf_fitz):
            fitz_page = pdf_fitz[page_idx]
            for img_info in fitz_page.get_images(full=True):
                try:
                    xref = img_info[0]
                    base_image = pdf_fitz.extract_image(xref)
                    img = Image.open(io.BytesIO(base_image["image"])).convert("RGB")
                    images.append(img)
                except Exception:
                    pass

        slides_data.append({
            "slide_number": page_idx + 1,
            "text": page_text.strip(),
            "images": images,
        })

    pdf_plumber.close()
    pdf_fitz.close()

    slides_data = _filter_slides_images(slides_data)

    num_pages = len(slides_data)
    avg_text = total_text_len / max(num_pages, 1)
    if avg_text < MIN_TEXT_PER_PAGE:
        return None

    return slides_data


def extract_lecture_notes_hybrid(
    slides_data: list[dict],
    *,
    provider: str = "openai",
    api_key: str | None = None,
    model: str | None = None,
    extra_instructions: str = "",
    on_progress: Any = None,
) -> dict[str, Any]:
    """Generate a notebook from pre-extracted text + images (hybrid path).

    Uses GPT-4o-mini when no images are present (fastest), GPT-4o when images exist.
    """
    api_key, model, _ = _resolve_provider(provider, api_key, model)

    total_images = sum(len(s["images"]) for s in slides_data)
    has_images = total_images > 0

    if has_images:
        use_model = model
    else:
        use_model = "gpt-4o-mini" if "openai" in provider or provider == "openai" else model

    if len(slides_data) <= HYBRID_CHUNK_SIZE:
        if on_progress:
            on_progress("analyzing", 1, 1)
        notebook, diagram_map = _hybrid_single_chunk(slides_data, api_key, use_model, provider, extra_instructions)
        return _embed_diagrams(notebook, diagram_map)

    return _hybrid_chunked(slides_data, api_key, use_model, model, provider, extra_instructions, on_progress)


def _resize_for_hybrid(img: Image.Image) -> Image.Image:
    """Resize an image to fit within HYBRID_IMAGE_MAX_PX, keeping aspect ratio."""
    w, h = img.size
    if max(w, h) <= HYBRID_IMAGE_MAX_PX:
        return img
    ratio = HYBRID_IMAGE_MAX_PX / max(w, h)
    return img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)


def _embed_diagrams(notebook: dict, diagram_map: dict[int, str]) -> dict:
    """Replace [DIAGRAM:N] placeholders with <img> tags in all content strings."""
    import re

    if not diagram_map:
        return notebook

    pattern = re.compile(r'\[DIAGRAM:(\d+)\]')

    def replace_in_str(s: str) -> str:
        def _sub(m):
            idx = int(m.group(1))
            b64 = diagram_map.get(idx)
            if b64:
                return f'<img src="data:image/png;base64,{b64}" class="nb-diagram" alt="Diagram {idx}" />'
            return m.group(0)
        return pattern.sub(_sub, s)

    for section in notebook.get("sections", []):
        if isinstance(section.get("content"), str):
            section["content"] = replace_in_str(section["content"])
        for sub in section.get("subsections", []):
            if isinstance(sub.get("content"), str):
                sub["content"] = replace_in_str(sub["content"])
            if isinstance(sub.get("bullets"), list):
                sub["bullets"] = [replace_in_str(b) if isinstance(b, str) else b for b in sub["bullets"]]

    return notebook


def _hybrid_single_chunk(
    slides_data: list[dict],
    api_key: str,
    model: str,
    provider: str,
    extra_instructions: str,
) -> tuple[dict[str, Any], dict[int, str]]:
    """Process a single chunk of hybrid slide data.

    Returns (notebook_dict, diagram_map) — placeholders are NOT yet replaced.
    """
    text_parts = []
    all_images = []

    for s in slides_data:
        header = f"--- Slide {s['slide_number']} ---"
        text_parts.append(header)
        if s["text"]:
            text_parts.append(s["text"])
        else:
            text_parts.append("[This slide contains only images/diagrams]")
        if s["images"]:
            text_parts.append(f"[{len(s['images'])} embedded image(s) attached below]")
            for img in s["images"]:
                all_images.append(_resize_for_hybrid(img))

    if len(all_images) > MAX_IMAGES_PER_CHUNK:
        all_images.sort(key=lambda im: im.size[0] * im.size[1], reverse=True)
        all_images = all_images[:MAX_IMAGES_PER_CHUNK]

    diagram_map: dict[int, str] = {}
    for i, img in enumerate(all_images):
        diagram_map[i + 1] = _image_to_base64(img, max_size=800)

    if all_images:
        user_text = (
            "Here is the extracted text content from university lecture slides.\n"
            f"There are {len(all_images)} embedded image(s) attached (numbered 1–{len(all_images)}).\n"
            "Use [DIAGRAM:N] placeholders in your content where each image should appear.\n"
            "Transform this into an intuitive, Socratic study guide.\n\n"
            + "\n".join(text_parts)
        )
    else:
        user_text = (
            "Here is the extracted text content from university lecture slides.\n"
            "Transform this into an intuitive, Socratic study guide.\n\n"
            + "\n".join(text_parts)
        )
    if extra_instructions:
        user_text += f"\n\nAdditional instructions: {extra_instructions}"

    try:
        if all_images:
            raw = _call_llm_raw(all_images, HYBRID_NOTEBOOK_PROMPT, user_text, api_key, model, provider, detail="low")
        else:
            raw = _call_text_llm(HYBRID_NOTEBOOK_PROMPT, user_text, api_key, model, provider)
    except Exception as e:
        err_str = str(e)
        if "429" in err_str or "rate_limit" in err_str:
            import time
            time.sleep(2)
            text_model = "gpt-4o-mini" if provider == "openai" else model
            raw = _call_text_llm(HYBRID_NOTEBOOK_PROMPT, user_text, api_key, text_model, provider)
            diagram_map = {}
        else:
            raise

    cleaned = _clean_json_response(raw)
    try:
        result = json.loads(cleaned)
    except json.JSONDecodeError as exc:
        raise ValueError(f"LLM returned invalid JSON. Raw:\n{raw[:1000]}") from exc

    return result, diagram_map


def _strip_diagram_placeholders(notebook: dict) -> dict:
    """Remove [DIAGRAM:N] placeholders from all content strings before merge."""
    import re
    pattern = re.compile(r'\s*\[DIAGRAM:\d+\]\s*')

    def clean(s):
        return pattern.sub(' ', s).strip() if isinstance(s, str) else s

    for section in notebook.get("sections", []):
        if isinstance(section.get("content"), str):
            section["content"] = clean(section["content"])
        for sub in section.get("subsections", []):
            if isinstance(sub.get("content"), str):
                sub["content"] = clean(sub["content"])
            if isinstance(sub.get("bullets"), list):
                sub["bullets"] = [clean(b) for b in sub["bullets"]]
    return notebook


MAX_TOTAL_DIAGRAMS = 5
MIN_SLOT_GAP = 3  # at least 3 subsections between any two diagrams


def _inject_diagrams_post_merge(
    notebook: dict,
    all_diagram_maps: list[dict[int, str]],
    num_chunks: int,
) -> dict:
    """Place diagrams into the merged notebook with even spacing.

    - Caps total diagrams at MAX_TOTAL_DIAGRAMS (keeps the largest).
    - Enforces MIN_SLOT_GAP subsections between any two diagrams.
    - Places each diagram at the end of a subsection's content.
    """
    all_images: list[tuple[int, str, int]] = []
    for chunk_idx, dmap in enumerate(all_diagram_maps):
        for key in sorted(dmap.keys()):
            b64 = dmap[key]
            all_images.append((chunk_idx, b64, len(b64)))

    if not all_images:
        return notebook

    if len(all_images) > MAX_TOTAL_DIAGRAMS:
        all_images.sort(key=lambda x: x[2], reverse=True)
        all_images = all_images[:MAX_TOTAL_DIAGRAMS]
        all_images.sort(key=lambda x: x[0])

    sections = notebook.get("sections", [])
    if not sections:
        return notebook

    slots: list[tuple[int, int]] = []
    for s_idx, section in enumerate(sections):
        subs = section.get("subsections", [])
        if subs:
            for sub_idx in range(len(subs)):
                slots.append((s_idx, sub_idx))
        else:
            slots.append((s_idx, -1))

    if not slots:
        return notebook

    n_imgs = len(all_images)
    spacing = max(MIN_SLOT_GAP, len(slots) // max(n_imgs, 1))
    chosen_slots = []
    for i in range(n_imgs):
        idx = min(i * spacing + spacing // 2, len(slots) - 1)
        while idx in chosen_slots and idx < len(slots) - 1:
            idx += 1
        chosen_slots.append(idx)

    for img_i, (chunk_idx, b64, _) in enumerate(all_images):
        slot_idx = chosen_slots[img_i]
        s_idx, sub_idx = slots[slot_idx]
        img_tag = f'<img src="data:image/png;base64,{b64}" class="nb-diagram" alt="Diagram" />'

        if sub_idx >= 0:
            sub = sections[s_idx]["subsections"][sub_idx]
            sub["content"] = sub.get("content", "") + "\n" + img_tag
        else:
            sections[s_idx]["content"] = sections[s_idx].get("content", "") + "\n" + img_tag

    return notebook


HYBRID_CHUNK_SIZE_TEXT = 15  # larger chunks when no images (text-only PDFs)

def _hybrid_chunked(
    slides_data: list[dict],
    api_key: str,
    chunk_model: str,
    merge_model: str,
    provider: str,
    extra_instructions: str,
    on_progress: Any = None,
) -> dict[str, Any]:
    """Process large hybrid slide data in parallel chunks, then merge.

    Dynamically picks chunk size: 15 for text-only PDFs, 5 for image-heavy.
    Uses hierarchical merge when >4 chunks to keep each merge prompt small.
    """
    import time

    total_images = sum(len(s["images"]) for s in slides_data)
    chunk_size = HYBRID_CHUNK_SIZE_TEXT if total_images == 0 else HYBRID_CHUNK_SIZE

    num_chunks = math.ceil(len(slides_data) / chunk_size)
    chunks = []
    for i in range(num_chunks):
        start = i * chunk_size
        end = min(start + chunk_size, len(slides_data))
        chunks.append(slides_data[start:end])

    if on_progress:
        on_progress("chunking", 0, num_chunks)

    partial_results: list[dict[str, Any]] = [None] * num_chunks
    chunk_diagram_maps: list[dict[int, str]] = [{}] * num_chunks
    completed = [0]

    print(f"  [hybrid] {num_chunks} chunks (size {chunk_size}, images={total_images}), max_workers={min(4, num_chunks)}")

    def process_chunk(idx: int, chunk: list[dict]) -> tuple[int, dict[str, Any], dict[int, str]]:
        time.sleep(0.2)
        t0 = time.time()
        n_imgs = sum(len(s["images"]) for s in chunk)
        chunk_instruction = (
            f"{extra_instructions}\n\n"
            f"These are slides {idx * chunk_size + 1}–"
            f"{idx * chunk_size + len(chunk)} "
            f"of {len(slides_data)} total slides from a single lecture."
        ).strip()
        result, dmap = _hybrid_single_chunk(chunk, api_key, chunk_model, provider, chunk_instruction)
        completed[0] += 1
        print(f"  [hybrid] chunk {idx+1}/{num_chunks} done in {time.time()-t0:.1f}s ({len(chunk)} slides, {n_imgs} imgs, {len(dmap)} diagrams)")
        if on_progress:
            on_progress("analyzing", completed[0], num_chunks)
        return idx, result, dmap

    max_workers = min(4, num_chunks)
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [executor.submit(process_chunk, i, c) for i, c in enumerate(chunks)]
        for future in concurrent.futures.as_completed(futures):
            idx, result, dmap = future.result()
            partial_results[idx] = _strip_diagram_placeholders(result)
            chunk_diagram_maps[idx] = dmap

    if on_progress:
        on_progress("merging", 0, 1)

    t_merge = time.time()
    merged = _hierarchical_merge(partial_results, api_key, merge_model, provider, on_progress)
    print(f"  [hybrid] merge step took {time.time()-t_merge:.1f}s")

    merged = _inject_diagrams_post_merge(merged, chunk_diagram_maps, num_chunks)
    total_diagrams = sum(len(d) for d in chunk_diagram_maps)
    if total_diagrams:
        print(f"  [hybrid] injected {total_diagrams} diagram(s) into merged notebook")

    if on_progress:
        on_progress("merging", 1, 1)

    return merged


def _hierarchical_merge(
    partials: list[dict[str, Any]],
    api_key: str,
    model: str,
    provider: str,
    on_progress: Any = None,
) -> dict[str, Any]:
    """Merge partials hierarchically: groups of 3-4, then merge intermediates.

    Keeps each LLM merge call small enough to fit in context.
    """
    MAX_GROUP = 4

    if len(partials) <= MAX_GROUP:
        return _merge_partial_notebooks(partials, api_key, model, provider)

    groups = []
    for i in range(0, len(partials), MAX_GROUP):
        groups.append(partials[i : i + MAX_GROUP])

    print(f"  [hybrid] hierarchical merge: {len(partials)} partials → {len(groups)} groups")
    intermediates = []
    for gi, group in enumerate(groups):
        print(f"  [hybrid] merging group {gi+1}/{len(groups)} ({len(group)} partials)")
        merged_group = _merge_partial_notebooks(group, api_key, model, provider)
        intermediates.append(merged_group)
        if on_progress:
            on_progress("merging", gi + 1, len(groups) + 1)

    print(f"  [hybrid] final merge of {len(intermediates)} intermediates")
    return _merge_partial_notebooks(intermediates, api_key, model, provider)


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def _image_to_base64(image: Image.Image, max_size: int = 2048) -> str:
    """Resize (if needed) and convert a PIL image to a base64 data-URI string."""
    w, h = image.size
    if max(w, h) > max_size:
        ratio = max_size / max(w, h)
        image = image.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)

    buf = io.BytesIO()
    image.save(buf, format="PNG")
    b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
    return b64


def _image_to_base64_for_detail(image: Image.Image, detail: str = "low") -> str:
    """Resize based on detail level and convert to base64.
    - 'low': resize to max 512px (fast, cheap — good for slides)
    - 'high': resize to max 2048px (detailed — good for handwriting/dense text)
    """
    max_size = 512 if detail == "low" else 2048
    return _image_to_base64(image, max_size=max_size)


def _clean_json_response(text: str) -> str:
    """Strip markdown fences or other wrapping from the LLM response."""
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
    return text.strip()


# ---------------------------------------------------------------------------
# Provider: OpenAI (GPT-4o)
# ---------------------------------------------------------------------------

def _extract_openai(
    images: list[Image.Image],
    api_key: str,
    model: str,
    extra_instructions: str,
) -> str:
    """Call OpenAI's vision API and return the raw text response."""
    from openai import OpenAI

    client = OpenAI(api_key=api_key)

    content: list[dict[str, Any]] = []
    user_text = "Here are the scanned pages of an exam paper. Extract all questions."
    if extra_instructions:
        user_text += f"\n\nAdditional instructions: {extra_instructions}"
    content.append({"type": "text", "text": user_text})

    for img in images:
        b64 = _image_to_base64(img)
        content.append(
            {
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "high"},
            }
        )

    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": content},
        ],
        temperature=0.1,
        max_tokens=16384,
    )
    return response.choices[0].message.content or ""


# ---------------------------------------------------------------------------
# Provider: Anthropic (Claude)
# ---------------------------------------------------------------------------

def _extract_anthropic(
    images: list[Image.Image],
    api_key: str,
    model: str,
    extra_instructions: str,
) -> str:
    """Call Anthropic's vision API and return the raw text response."""
    try:
        import anthropic
    except ImportError:
        raise ImportError(
            "anthropic package is required for Claude provider. "
            "Install with: pip install anthropic"
        )

    client = anthropic.Anthropic(api_key=api_key)

    content: list[dict[str, Any]] = []

    # Add images first for Claude
    for img in images:
        b64 = _image_to_base64(img)
        content.append(
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/png",
                    "data": b64,
                },
            }
        )

    user_text = "Here are the scanned pages of an exam paper. Extract all questions."
    if extra_instructions:
        user_text += f"\n\nAdditional instructions: {extra_instructions}"
    content.append({"type": "text", "text": user_text})

    response = client.messages.create(
        model=model,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": content}],
        temperature=0.1,
        max_tokens=16384,
    )
    return response.content[0].text


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

# Map of provider name → (default model, env var for key, extractor function)
PROVIDERS = {
    "openai": {
        "default_model": "gpt-4o",
        "env_var": "OPENAI_API_KEY",
        "extractor": _extract_openai,
    },
    "anthropic": {
        "default_model": "claude-sonnet-4-20250514",
        "env_var": "ANTHROPIC_API_KEY",
        "extractor": _extract_anthropic,
    },
    "kimi": {
        "default_model": "moonshotai/kimi-k2.5",
        "env_var": "KIMI_API_KEY",
        "base_url": "https://integrate.api.nvidia.com/v1",
    },
}

# Kimi uses the OpenAI-compatible API, so it shares the same calling convention.
# We handle it inside _call_llm_raw and _call_text_llm by switching base_url.


def _resolve_provider(
    provider: str,
    api_key: str | None = None,
    model: str | None = None,
) -> tuple[str, str, Any]:
    """Resolve provider config, API key, and model. Returns (api_key, model, extractor_fn)."""
    provider = provider.lower()
    if provider not in PROVIDERS:
        raise ValueError(
            f"Unknown provider '{provider}'. Choose from: {', '.join(PROVIDERS.keys())}"
        )

    prov_config = PROVIDERS[provider]
    api_key = api_key or os.getenv(prov_config["env_var"])
    if not api_key:
        raise ValueError(
            f"API key required for {provider}. "
            f"Set {prov_config['env_var']} env var or pass --api-key."
        )

    model = model or os.getenv(
        f"{provider.upper()}_MODEL",
        prov_config["default_model"],
    )

    return api_key, model, prov_config.get("extractor")


def _call_vision_llm(
    images: list[Image.Image],
    system_prompt: str,
    user_text: str,
    *,
    provider: str = "openai",
    api_key: str | None = None,
    model: str | None = None,
    detail: str = "high",
) -> dict[str, Any]:
    """Generic helper: send images + prompt to vision LLM and parse the JSON response."""
    api_key, model, extractor_fn = _resolve_provider(provider, api_key, model)

    raw = _call_llm_raw(images, system_prompt, user_text, api_key, model, provider, detail=detail)
    cleaned = _clean_json_response(raw)

    try:
        result = json.loads(cleaned)
    except json.JSONDecodeError as exc:
        raise ValueError(
            f"LLM returned invalid JSON. Raw response:\n{raw[:1000]}"
        ) from exc

    return result


def _call_llm_raw(
    images: list[Image.Image],
    system_prompt: str,
    user_text: str,
    api_key: str,
    model: str,
    provider: str,
    detail: str = "high",
) -> str:
    """Low-level: send images to a vision LLM with a custom system prompt.
    
    Args:
        detail: 'low' for slides (512px, 85 tokens/img) or 'high' for dense docs (2048px).
    """
    if provider in ("openai", "kimi"):
        from openai import OpenAI

        client_kwargs = {"api_key": api_key}
        if provider == "kimi":
            client_kwargs["base_url"] = PROVIDERS["kimi"]["base_url"]

        client = OpenAI(**client_kwargs)
        content: list[dict[str, Any]] = [{"type": "text", "text": user_text}]
        for img in images:
            b64 = _image_to_base64_for_detail(img, detail=detail)
            content.append(
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{b64}", "detail": detail},
                }
            )
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": content},
            ],
            temperature=0.2,
            max_tokens=16384,
        )
        return response.choices[0].message.content or ""

    elif provider == "anthropic":
        try:
            import anthropic
        except ImportError:
            raise ImportError(
                "anthropic package is required for Claude provider. "
                "Install with: pip install anthropic"
            )

        client = anthropic.Anthropic(api_key=api_key)
        content_parts: list[dict[str, Any]] = []
        for img in images:
            b64 = _image_to_base64_for_detail(img, detail=detail)
            content_parts.append(
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": "image/png",
                        "data": b64,
                    },
                }
            )
        content_parts.append({"type": "text", "text": user_text})

        response = client.messages.create(
            model=model,
            system=system_prompt,
            messages=[{"role": "user", "content": content_parts}],
            temperature=0.2,
            max_tokens=16384,
        )
        return response.content[0].text
    else:
        raise ValueError(f"Unknown provider: {provider}")


def _call_text_llm(
    system_prompt: str,
    user_text: str,
    api_key: str,
    model: str,
    provider: str,
    max_tokens: int = 16384,
) -> str:
    """Call an LLM with text-only input (no images). Used for question matching and merging."""
    if provider in ("openai", "kimi"):
        from openai import OpenAI

        client_kwargs = {"api_key": api_key}
        if provider == "kimi":
            client_kwargs["base_url"] = PROVIDERS["kimi"]["base_url"]

        client = OpenAI(**client_kwargs)
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_text},
            ],
            temperature=0.1,
            max_tokens=max_tokens,
        )
        return response.choices[0].message.content or ""

    elif provider == "anthropic":
        try:
            import anthropic
        except ImportError:
            raise ImportError("anthropic package is required. Install with: pip install anthropic")

        client = anthropic.Anthropic(api_key=api_key)
        response = client.messages.create(
            model=model,
            system=system_prompt,
            messages=[{"role": "user", "content": user_text}],
            temperature=0.1,
            max_tokens=max_tokens,
        )
        return response.content[0].text
    else:
        raise ValueError(f"Unknown provider: {provider}")


# ---------------------------------------------------------------------------
# Public API – Exam Paper Extraction (original)
# ---------------------------------------------------------------------------

def extract_exam_paper(
    images: list[Image.Image],
    *,
    provider: str = "openai",
    api_key: str | None = None,
    model: str | None = None,
    extra_instructions: str = "",
    detail: str = "high",
) -> dict[str, Any]:
    """
    Send page images to the vision LLM and return the structured exam JSON.

    Args:
        images: List of PIL images (one per page).
        provider: 'openai', 'anthropic', or 'kimi'.
        api_key: API key. Falls back to env var based on provider.
        model: Model name. Falls back to provider default.
        extra_instructions: Additional context to append to the system prompt.
        detail: 'high' (default for exams with fine print) or 'low'.

    Returns:
        Parsed JSON dict conforming to the ExamPaper schema.
    """
    api_key, model, _ = _resolve_provider(provider, api_key, model)

    user_text = "Here are the scanned pages of an exam paper. Extract all questions."
    if extra_instructions:
        user_text += f"\n\nAdditional instructions: {extra_instructions}"

    raw = _call_llm_raw(images, SYSTEM_PROMPT, user_text, api_key, model, provider, detail=detail)
    cleaned = _clean_json_response(raw)

    try:
        result = json.loads(cleaned)
    except json.JSONDecodeError as exc:
        raise ValueError(
            f"LLM returned invalid JSON. Raw response:\n{raw[:1000]}"
        ) from exc

    return result


# ---------------------------------------------------------------------------
# Public API – Notebook / Study Guide Generation
# ---------------------------------------------------------------------------

CHUNK_SIZE = 8  # slides per chunk — sweet spot for speed vs context


def extract_lecture_notes(
    images: list[Image.Image],
    *,
    provider: str = "openai",
    api_key: str | None = None,
    model: str | None = None,
    extra_instructions: str = "",
    detail: str = "low",
    on_progress: Any = None,
) -> dict[str, Any]:
    """
    Smart lecture note extraction — automatically chunks large slide decks.

    For <= CHUNK_SIZE slides: single call (fast).
    For > CHUNK_SIZE slides: parallel chunked processing + merge.

    Args:
        images: List of PIL images (one per page/slide).
        provider: 'openai', 'anthropic', or 'kimi'.
        api_key: API key. Falls back to env var based on provider.
        model: Model name. Falls back to provider default.
        extra_instructions: Additional context.
        detail: 'low' (default, good for slides) or 'high' (for dense handwriting).
        on_progress: Optional callback(stage: str, current: int, total: int) for progress.

    Returns:
        Parsed JSON dict conforming to the Notebook schema.
    """
    api_key, model, _ = _resolve_provider(provider, api_key, model)

    if len(images) <= CHUNK_SIZE:
        # Small deck — single call, no chunking needed
        if on_progress:
            on_progress("analyzing", 1, 1)
        return _extract_single_chunk(images, api_key, model, provider, extra_instructions, detail)

    # Large deck — chunk + parallel + merge
    return _extract_chunked(images, api_key, model, provider, extra_instructions, detail, on_progress)


def _extract_single_chunk(
    images: list[Image.Image],
    api_key: str,
    model: str,
    provider: str,
    extra_instructions: str,
    detail: str,
) -> dict[str, Any]:
    """Process a single batch of images into a notebook."""
    user_text = (
        "Here are the pages/slides from a university lecture. "
        "Transform this content into an intuitive, Socratic study guide."
    )
    if extra_instructions:
        user_text += f"\n\nAdditional instructions: {extra_instructions}"

    raw = _call_llm_raw(images, NOTEBOOK_SYSTEM_PROMPT, user_text, api_key, model, provider, detail=detail)
    cleaned = _clean_json_response(raw)

    try:
        result = json.loads(cleaned)
    except json.JSONDecodeError as exc:
        raise ValueError(
            f"LLM returned invalid JSON. Raw response:\n{raw[:1000]}"
        ) from exc

    return result


def _extract_chunked(
    images: list[Image.Image],
    api_key: str,
    model: str,
    provider: str,
    extra_instructions: str,
    detail: str,
    on_progress: Any = None,
) -> dict[str, Any]:
    """Process a large slide deck in parallel chunks, then merge into one notebook."""
    # Split into chunks
    num_chunks = math.ceil(len(images) / CHUNK_SIZE)
    chunks = []
    for i in range(num_chunks):
        start = i * CHUNK_SIZE
        end = min(start + CHUNK_SIZE, len(images))
        chunks.append(images[start:end])

    if on_progress:
        on_progress("chunking", 0, num_chunks)

    # Process chunks in parallel (max 4 concurrent to avoid rate limits)
    partial_results: list[dict[str, Any]] = [None] * num_chunks  # type: ignore
    completed = [0]

    def process_chunk(idx: int, chunk_images: list[Image.Image]) -> tuple[int, dict[str, Any]]:
        chunk_instruction = (
            f"{extra_instructions}\n\n"
            f"These are slides {idx * CHUNK_SIZE + 1}–{idx * CHUNK_SIZE + len(chunk_images)} "
            f"of {len(images)} total slides from a single lecture."
        ).strip()

        result = _extract_single_chunk(chunk_images, api_key, model, provider, chunk_instruction, detail)

        completed[0] += 1
        if on_progress:
            on_progress("analyzing", completed[0], num_chunks)

        return idx, result

    max_workers = min(4, num_chunks)
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [
            executor.submit(process_chunk, i, chunk)
            for i, chunk in enumerate(chunks)
        ]
        for future in concurrent.futures.as_completed(futures):
            idx, result = future.result()
            partial_results[idx] = result

    # Merge all partial results into one cohesive notebook
    if on_progress:
        on_progress("merging", 0, 1)

    merged = _merge_partial_notebooks(partial_results, api_key, model, provider)

    if on_progress:
        on_progress("merging", 1, 1)

    return merged


def _repair_truncated_json(text: str) -> str:
    """Attempt to repair JSON truncated mid-output by closing open brackets/braces."""
    stack = []
    in_string = False
    escape = False
    for ch in text:
        if escape:
            escape = False
            continue
        if ch == '\\' and in_string:
            escape = True
            continue
        if ch == '"' and not escape:
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch in ('{', '['):
            stack.append('}' if ch == '{' else ']')
        elif ch in ('}', ']'):
            if stack and stack[-1] == ch:
                stack.pop()

    if in_string:
        text += '"'
    text = text.rstrip().rstrip(',')
    while stack:
        text += stack.pop()
    return text


def _merge_partial_notebooks(
    partials: list[dict[str, Any]],
    api_key: str,
    model: str,
    provider: str,
) -> dict[str, Any]:
    """Use the LLM to merge multiple partial notebooks into one cohesive guide."""
    if len(partials) == 1:
        return partials[0]

    chunks_text = []
    for i, p in enumerate(partials):
        compact = json.dumps(p, separators=(',', ':'))
        chunks_text.append(f"=== CHUNK {i + 1} of {len(partials)} ===\n{compact}")

    user_text = (
        "Here are partial study guides generated from different batches of the same lecture.\n"
        "Merge them into ONE cohesive, comprehensive study guide.\n\n"
        + "\n\n".join(chunks_text)
    )

    raw = _call_text_llm(MERGE_SYSTEM_PROMPT, user_text, api_key, model, provider, max_tokens=32768)
    cleaned = _clean_json_response(raw)

    try:
        merged = json.loads(cleaned)
    except json.JSONDecodeError:
        try:
            repaired = _repair_truncated_json(cleaned)
            merged = json.loads(repaired)
            print("  [merge] repaired truncated JSON successfully")
        except (json.JSONDecodeError, Exception):
            print("  [merge] LLM JSON failed, using manual merge fallback")
            merged = _manual_merge_fallback(partials)

    return merged


def _manual_merge_fallback(partials: list[dict[str, Any]]) -> dict[str, Any]:
    """Simple concatenation fallback if the LLM merge fails."""
    base = partials[0].copy()
    all_sections = []
    all_chat = []
    all_tags = set()

    for p in partials:
        for section in p.get("sections", []):
            all_sections.append(section)
            for tag in section.get("tags", []):
                all_tags.add(tag)
        cr = p.get("chatResponses", [])
        if isinstance(cr, list):
            all_chat.extend(cr)
        elif isinstance(cr, dict):
            for k, v in cr.items():
                all_chat.append({"keywords": k, "response": v})

    base["sections"] = all_sections
    base["chatResponses"] = all_chat
    return base


def match_questions_to_notebook(
    notebook: dict[str, Any],
    papers: list[dict[str, Any]],
    *,
    provider: str = "openai",
    api_key: str | None = None,
    model: str | None = None,
) -> list[str]:
    """
    Use the LLM to match past paper questions to notebook sections based on topic overlap.

    Args:
        notebook: The generated notebook dict (with sections and tags).
        papers: List of exam paper dicts (each with a 'questions' list).
        provider: 'openai' or 'anthropic'.
        api_key: API key.
        model: Model name.

    Returns:
        List of matched question IDs.
    """
    api_key, model, _ = _resolve_provider(provider, api_key, model)

    # Build rich context for the LLM — include section content, not just tags
    sections_context = []
    for section in notebook.get("sections", []):
        sec_info = {
            "title": section.get("title", ""),
            "tags": section.get("tags", []),
            "content_preview": section.get("content", "")[:200],
        }
        # Include subsection titles for additional context
        subs = section.get("subsections") or []
        sec_info["subtopics"] = [s.get("title", "") for s in subs]
        sections_context.append(sec_info)

    all_questions = []
    for paper in papers:
        paper_id = paper.get("id", "unknown")
        for q in paper.get("questions", []):
            all_questions.append({
                "id": f"{paper_id}_{q['id']}",
                "original_id": q["id"],
                "paper_id": paper_id,
                "paper_title": paper.get("title", ""),
                "type": q.get("type", "unknown"),
                "text": q.get("text", ""),
                "equation": q.get("equation"),
            })

    user_text = (
        f"### Notebook Sections (with topics and content)\n"
        f"{json.dumps(sections_context, indent=2)}\n\n"
        f"### Available Past Paper Questions\n"
        f"{json.dumps(all_questions, indent=2)}\n\n"
        "Match the most relevant questions to these notebook topics. "
        "Use semantic understanding — tags may use different vocabulary than questions. "
        "Return their IDs (the compound id field like 'paper_id_qN')."
    )

    raw = _call_text_llm(MATCHING_SYSTEM_PROMPT, user_text, api_key, model, provider)
    cleaned = _clean_json_response(raw)

    try:
        matched_ids = json.loads(cleaned)
    except json.JSONDecodeError:
        # Fallback: basic keyword matching if LLM fails
        matched_ids = _fallback_keyword_match(all_tags, all_questions)

    return matched_ids


def _fallback_keyword_match(
    tags: list[str],
    questions: list[dict[str, Any]],
    max_matches: int = 8,
) -> list[str]:
    """Simple keyword-based fallback matcher if the LLM call fails."""
    tag_set = set(t.lower() for t in tags)
    scored = []
    for q in questions:
        text_lower = q["text"].lower()
        score = sum(1 for tag in tag_set if tag in text_lower)
        if score > 0:
            scored.append((score, q["id"]))
    scored.sort(key=lambda x: x[0], reverse=True)
    return [qid for _, qid in scored[:max_matches]]
