"""FastAPI server for Coast — auth, notebooks, sessions, and OCR pipeline."""

from __future__ import annotations

import json
import math
import os
import shutil
import tempfile
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv
from fastapi import Depends, FastAPI, File, Form, Header, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel

load_dotenv()

from auth import create_access_token, decode_access_token, hash_password, verify_password
from database import (
    ActivityEvent,
    ChatMessage,
    FolderSource,
    LessonNotes,
    Paper,
    QuizSession,
    SavedNotebook,
    SessionAnswer,
    SessionLocal,
    SkillProfile,
    SourceImage,
    StudyFolder,
    TutorMemo,
    User,
    UserFeedback,
    init_db,
    load_papers_from_json,
)
import threading
import tutor
import spaced_rep
import rag
import lesson


def _bg_post_process(user_id: int, folder: str, nb_id: str, notebook_data: dict):
    """Run RAG embedding and card creation in background after notebook save."""
    if folder:
        try:
            rag.embed_notebook(user_id, folder, nb_id, notebook_data)
        except Exception:
            import traceback
            traceback.print_exc()
    try:
        spaced_rep.create_cards_for_notebook(user_id, nb_id, notebook_data)
    except Exception:
        import traceback
        traceback.print_exc()
    print(f"[BG] Post-processing done for notebook {nb_id}")
from viz_router import viz_router

app = FastAPI(title="Coast API", version="2.0.0")
app.include_router(viz_router)

_ALLOWED_ORIGINS = [
    "http://localhost:5173", "http://localhost:5174", "http://localhost:3000",
    "https://dist-delta-eight-99.vercel.app",
]
_extra_origin = os.getenv("FRONTEND_URL", "")
if _extra_origin:
    _ALLOWED_ORIGINS.append(_extra_origin.rstrip("/"))

app.add_middleware(
    CORSMiddleware,
    allow_origins=_ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"],
    max_age=600,
)

from starlette.requests import Request as StarletteRequest
from starlette.responses import Response as StarletteResponse

_PROD_ORIGIN = "https://dist-delta-eight-99.vercel.app"

@app.options("/{rest:path}")
async def preflight_catchall(request: StarletteRequest):
    """Catch-all for OPTIONS preflight — always return CORS headers."""
    origin = request.headers.get("origin", _PROD_ORIGIN)
    allowed = origin if origin in _ALLOWED_ORIGINS else _PROD_ORIGIN
    return StarletteResponse(
        status_code=204,
        headers={
            "Access-Control-Allow-Origin": allowed,
            "Access-Control-Allow-Credentials": "true",
            "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS, PATCH",
            "Access-Control-Allow-Headers": "Authorization, Content-Type, X-Requested-With, Accept, Origin",
            "Access-Control-Max-Age": "600",
        },
    )

@app.middleware("http")
async def cors_safety_net(request: StarletteRequest, call_next):
    """Last-resort: inject CORS headers if CORSMiddleware didn't."""
    response = await call_next(request)
    if "access-control-allow-origin" not in response.headers:
        origin = request.headers.get("origin", _PROD_ORIGIN)
        allowed = origin if origin in _ALLOWED_ORIGINS else _PROD_ORIGIN
        response.headers["Access-Control-Allow-Origin"] = allowed
        response.headers["Access-Control-Allow-Credentials"] = "true"
    return response

_cors_headers = {
    "Access-Control-Allow-Origin": _PROD_ORIGIN,
    "Access-Control-Allow-Credentials": "true",
    "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS, PATCH",
    "Access-Control-Allow-Headers": "Authorization, Content-Type",
}

@app.exception_handler(HTTPException)
async def http_exception_handler(request: StarletteRequest, exc: HTTPException):
    return JSONResponse(
        status_code=exc.status_code,
        content={"detail": exc.detail},
        headers=_cors_headers,
    )

@app.exception_handler(Exception)
async def global_exception_handler(request: StarletteRequest, exc: Exception):
    import traceback
    traceback.print_exc()
    return JSONResponse(
        status_code=500,
        content={"detail": "Internal server error"},
        headers=_cors_headers,
    )

GENERATED_DIR = Path(__file__).parent / "generated"
GENERATED_DIR.mkdir(exist_ok=True)

_PERSISTENT_DISK = Path("/data")
_default_uploads = str(_PERSISTENT_DISK / "folder_uploads") if _PERSISTENT_DISK.is_dir() else str(Path(__file__).parent / "folder_uploads")
FOLDER_UPLOADS_DIR = Path(os.environ.get("FOLDER_UPLOADS_DIR", _default_uploads))
FOLDER_UPLOADS_DIR.mkdir(parents=True, exist_ok=True)

SOURCE_IMAGES_DIR = FOLDER_UPLOADS_DIR / "images"
SOURCE_IMAGES_DIR.mkdir(parents=True, exist_ok=True)

PAPERS_DIR = Path(os.getenv("PAPERS_DIR", str(Path(__file__).parent.parent / "Coast" / "testing" / "src" / "data")))

from curated_config import curated_source_uid as _curated_uid, get_lesson_structure

QUESTIONS_PER_BATCH = 10

RATE_LIMIT_CHAT_MESSAGES = 500      # max Pedro messages per user per week
RATE_LIMIT_NOTEBOOKS = 20           # max notebook uploads per user
RATE_LIMIT_WINDOW_DAYS = 7          # rolling window for chat limit

_live_users: dict[int, dict] = {}
HEARTBEAT_TIMEOUT = 60


# ═══════════════════════════════════════════════════════════════════════════
# STARTUP
# ═══════════════════════════════════════════════════════════════════════════

@app.on_event("startup")
def on_startup():
    init_db()
    if PAPERS_DIR.exists():
        load_papers_from_json(PAPERS_DIR)
        print(f"Loaded papers from {PAPERS_DIR}")
    from paper_scanner import load_scanned_into_db
    load_scanned_into_db()

    import threading
    def _bg_curated():
        try:
            from curated_config import ingest_curated_sources
            ingest_curated_sources()
        except Exception:
            import traceback; traceback.print_exc()
    threading.Thread(target=_bg_curated, daemon=True).start()


# ═══════════════════════════════════════════════════════════════════════════
# AUTH DEPENDENCY
# ═══════════════════════════════════════════════════════════════════════════

def get_current_user(authorization: Optional[str] = Header(None)) -> User:
    """Extract and verify the JWT from the Authorization header."""
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(401, "Not authenticated")

    token = authorization.split(" ", 1)[1]
    payload = decode_access_token(token)
    if not payload:
        raise HTTPException(401, "Invalid or expired token")

    db = SessionLocal()
    user = db.query(User).filter(User.id == int(payload["sub"])).first()
    db.close()

    if not user:
        raise HTTPException(401, "User not found")
    return user


ADMIN_EMAIL = os.getenv("ADMIN_EMAIL", "andreaf.fraschetti@gmail.com")


def _get_user_usage(user_id: int):
    """Return current usage counts for rate limiting. Admin gets unlimited."""
    from datetime import timedelta

    db = SessionLocal()
    try:
        admin = db.query(User).filter(User.id == user_id).first()
        if admin and admin.email == ADMIN_EMAIL:
            return {
                "chat_messages_used": 0,
                "chat_messages_limit": 999999,
                "chat_messages_remaining": 999999,
                "notebooks_used": 0,
                "notebooks_limit": 999999,
                "notebooks_remaining": 999999,
            }

        cutoff = datetime.now(timezone.utc) - timedelta(days=RATE_LIMIT_WINDOW_DAYS)
        chat_count = (
            db.query(ChatMessage)
            .filter(
                ChatMessage.user_id == user_id,
                ChatMessage.role == "user",
                ChatMessage.created_at >= cutoff,
            )
            .count()
        )
        notebook_count = (
            db.query(SavedNotebook)
            .filter(
                SavedNotebook.user_id == user_id,
                SavedNotebook.is_premade == False,
                SavedNotebook.deleted_at == None,
            )
            .count()
        )
        return {
            "chat_messages_used": chat_count,
            "chat_messages_limit": RATE_LIMIT_CHAT_MESSAGES,
            "chat_messages_remaining": max(0, RATE_LIMIT_CHAT_MESSAGES - chat_count),
            "notebooks_used": notebook_count,
            "notebooks_limit": RATE_LIMIT_NOTEBOOKS,
            "notebooks_remaining": max(0, RATE_LIMIT_NOTEBOOKS - notebook_count),
        }
    finally:
        db.close()


@app.get("/api/usage")
def get_usage(user: User = Depends(get_current_user)):
    """Return the user's current usage against rate limits."""
    return _get_user_usage(user.id)


# ═══════════════════════════════════════════════════════════════════════════
# AUTH ENDPOINTS
# ═══════════════════════════════════════════════════════════════════════════

class RegisterRequest(BaseModel):
    email: str
    name: str
    password: str
    course: str = ""


class LoginRequest(BaseModel):
    email: str
    password: str


@app.post("/api/auth/register")
def register(req: RegisterRequest):
    db = SessionLocal()
    try:
        existing = db.query(User).filter(User.email == req.email.lower().strip()).first()
        if existing:
            raise HTTPException(400, "Email already registered")

        user = User(
            email=req.email.lower().strip(),
            name=req.name.strip(),
            password_hash=hash_password(req.password),
            course=req.course.strip(),
        )
        db.add(user)
        db.commit()
        db.refresh(user)

        token = create_access_token(user.id, user.email)
        return {
            "token": token,
            "user": {"id": user.id, "email": user.email, "name": user.name, "course": user.course, "onboarding_completed": False},
        }
    finally:
        db.close()


@app.post("/api/auth/login")
def login(req: LoginRequest):
    db = SessionLocal()
    try:
        user = db.query(User).filter(User.email == req.email.lower().strip()).first()
        if not user or not verify_password(req.password, user.password_hash):
            raise HTTPException(401, "Invalid email or password")

        token = create_access_token(user.id, user.email)
        return {
            "token": token,
            "user": {"id": user.id, "email": user.email, "name": user.name, "course": user.course, "onboarding_completed": bool(user.onboarding_completed)},
        }
    finally:
        db.close()


@app.get("/api/auth/me")
def get_me(user: User = Depends(get_current_user)):
    return {"id": user.id, "email": user.email, "name": user.name, "course": user.course, "onboarding_completed": bool(user.onboarding_completed)}


class OnboardingRequest(BaseModel):
    preferences: dict


@app.post("/api/auth/onboarding")
def complete_onboarding(req: OnboardingRequest, user: User = Depends(get_current_user)):
    db = SessionLocal()
    try:
        db_user = db.query(User).filter(User.id == user.id).first()
        if not db_user:
            raise HTTPException(404, "User not found")
        db_user.learning_preferences = json.dumps(req.preferences)
        db_user.onboarding_completed = True
        db.commit()
        return {
            "id": db_user.id,
            "email": db_user.email,
            "name": db_user.name,
            "course": db_user.course,
            "onboarding_completed": True,
        }
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# PAPERS / QUESTIONS
# ═══════════════════════════════════════════════════════════════════════════

@app.get("/api/papers")
def list_papers():
    """List all available papers."""
    db = SessionLocal()
    try:
        papers = db.query(Paper).all()
        return [
            {
                "id": p.paper_id,
                "title": p.title,
                "description": p.description,
                "course": p.course,
                "questionCount": p.question_count,
            }
            for p in papers
        ]
    finally:
        db.close()


@app.get("/api/papers/{paper_id}")
def get_paper(paper_id: str):
    """Get a full paper with all questions."""
    db = SessionLocal()
    try:
        paper = db.query(Paper).filter(Paper.paper_id == paper_id).first()
        if not paper:
            raise HTTPException(404, "Paper not found")
        return {
            "id": paper.paper_id,
            "title": paper.title,
            "description": paper.description,
            "course": paper.course,
            "questions": json.loads(paper.questions_json),
        }
    finally:
        db.close()


@app.get("/api/papers/{paper_id}/questions")
def get_questions_paginated(paper_id: str, batch: int = 1, user: User = Depends(get_current_user)):
    """Get a batch of questions (10 at a time). Tracks what the user has already seen."""
    db = SessionLocal()
    try:
        paper = db.query(Paper).filter(Paper.paper_id == paper_id).first()
        if not paper:
            raise HTTPException(404, "Paper not found")

        all_questions = json.loads(paper.questions_json)

        # Get questions already answered by this user for this paper
        answered_ids = set()
        prev_sessions = (
            db.query(QuizSession)
            .filter(QuizSession.user_id == user.id, QuizSession.paper_id == paper_id, QuizSession.completed == True)
            .all()
        )
        for sess in prev_sessions:
            for ans in sess.answers:
                answered_ids.add(ans.question_id)

        # Filter to unanswered questions first, then pad with answered if needed
        unanswered = [q for q in all_questions if q["id"] not in answered_ids]
        remaining = unanswered if unanswered else all_questions

        # Paginate
        start = (batch - 1) * QUESTIONS_PER_BATCH
        end = start + QUESTIONS_PER_BATCH
        batch_questions = remaining[start:end]
        total_batches = math.ceil(len(remaining) / QUESTIONS_PER_BATCH)

        return {
            "paper_id": paper_id,
            "paper_title": paper.title,
            "batch": batch,
            "total_batches": total_batches,
            "total_questions": len(remaining),
            "questions": batch_questions,
            "has_more": end < len(remaining),
        }
    finally:
        db.close()


@app.get("/api/questions/by-tags")
def get_questions_by_tags(tags: str, batch: int = 1, course: str = ""):
    """Get questions matching any of the given tags (comma-separated). Optionally filter by course."""
    tag_list = [t.strip().lower() for t in tags.split(",") if t.strip()]
    if not tag_list:
        raise HTTPException(400, "No tags provided")

    db = SessionLocal()
    try:
        if course:
            papers = db.query(Paper).filter(Paper.course == course).all()
        else:
            papers = db.query(Paper).all()
        matched = []

        for paper in papers:
            questions = json.loads(paper.questions_json)
            for q in questions:
                q_text = q.get("text", "").lower()
                q_tags = [t.lower() for t in (q.get("tags", []) or [])]
                q_key_terms = [t.lower() for t in (q.get("keyTerms", []) or [])]
                searchable = q_text + " " + " ".join(q_tags) + " ".join(q_key_terms)

                if any(tag in searchable for tag in tag_list):
                    matched.append({**q, "_paper_id": paper.paper_id, "_paper_title": paper.title})

        # Paginate
        start = (batch - 1) * QUESTIONS_PER_BATCH
        end = start + QUESTIONS_PER_BATCH
        total_batches = math.ceil(len(matched) / QUESTIONS_PER_BATCH) if matched else 1

        return {
            "tags": tag_list,
            "batch": batch,
            "total_batches": total_batches,
            "total_matched": len(matched),
            "questions": matched[start:end],
            "has_more": end < len(matched),
        }
    finally:
        db.close()


@app.get("/api/folders/{folder_name}/section-questions")
def get_section_questions(
    folder_name: str,
    user: User = Depends(get_current_user),
    max_questions: int = 5,
):
    """Find past paper questions matching the current lesson section's topics.

    Scores each question by overlap between question tags and section key_topics.
    Returns empty list if no matching questions — frontend hides the exam block.
    """
    from database import CourseOutline
    db = SessionLocal()
    try:
        outline = db.query(CourseOutline).filter(
            CourseOutline.user_id == user.id,
            CourseOutline.folder_name == folder_name,
        ).first()
        if not outline:
            return {"questions": [], "section_title": ""}

        sections = json.loads(outline.outline_json)
        idx = outline.current_section
        if idx >= len(sections):
            return {"questions": [], "section_title": ""}

        section = sections[idx]
        section_title = section.get("title", "")
        key_topics = [t.lower() for t in section.get("key_topics", [])]
        objectives = [o.lower() for o in section.get("learning_objectives", [])]

        search_terms = key_topics + objectives + [section_title.lower()]
        search_words = set()
        for term in search_terms:
            for word in term.split():
                if len(word) > 3:
                    search_words.add(word)

        from curated_config import get_course_for_folder
        course_code = get_course_for_folder(folder_name)
        if not course_code:
            return {"questions": [], "section_title": section_title}
        papers = db.query(Paper).filter(Paper.course == course_code).all()

        answered_ids = set()
        answered_rows = (
            db.query(SessionAnswer.question_id)
            .join(QuizSession)
            .filter(QuizSession.user_id == user.id)
            .all()
        )
        for row in answered_rows:
            answered_ids.add(row[0])

        scored = []
        for paper in papers:
            questions = json.loads(paper.questions_json)
            for q in questions:
                compound_id = f"{paper.paper_id}_{q.get('id', '')}"
                if compound_id in answered_ids or q.get("id", "") in answered_ids:
                    continue

                q_tags = [t.lower() for t in (q.get("tags", []) or [])]
                q_key_terms = [t.lower() for t in (q.get("keyTerms", []) or [])]
                q_text_lower = q.get("text", "").lower()

                score = 0
                for topic in key_topics:
                    for tag in q_tags:
                        if topic in tag or tag in topic:
                            score += 10
                    for kt in q_key_terms:
                        if topic in kt or kt in topic:
                            score += 5
                    if topic in q_text_lower:
                        score += 3
                for word in search_words:
                    if any(word in tag for tag in q_tags):
                        score += 2
                    if word in q_text_lower:
                        score += 1

                if score > 0:
                    scored.append((q, paper, score))

        scored.sort(key=lambda x: -x[2])
        top = scored[:max_questions]

        result_questions = []
        for q, paper, sc in top:
            out = {**q}
            out["_paper_id"] = paper.paper_id
            out["_paper_title"] = paper.title
            out["_match_score"] = sc
            result_questions.append(out)

        return {
            "questions": result_questions,
            "section_title": section_title,
            "section_topics": key_topics,
        }
    finally:
        db.close()


@app.post("/api/folders/{folder_name}/section-answers")
def submit_section_answers(
    folder_name: str,
    request: dict = {},
    user: User = Depends(get_current_user),
):
    """Submit answers for section past-paper questions and track per-topic results."""
    answers = request.get("answers", [])
    if not answers:
        return {"status": "no_answers"}

    db = SessionLocal()
    try:
        session = QuizSession(
            user_id=user.id,
            paper_id=f"lesson_{folder_name}",
            paper_title=f"Lesson: {folder_name}",
            score=sum(1 for a in answers if a.get("is_correct")),
            total=len(answers),
            completed=True,
            completed_at=datetime.now(timezone.utc),
        )
        db.add(session)
        db.flush()

        for a in answers:
            sa = SessionAnswer(
                session_id=session.id,
                question_id=a.get("question_id", ""),
                question_text=a.get("question_text", ""),
                user_answer=a.get("user_answer", ""),
                correct_answer=a.get("correct_answer", ""),
                is_correct=a.get("is_correct", False),
                tags_json=json.dumps(a.get("tags", [])),
            )
            db.add(sa)

        db.commit()

        def _bg_update_skills():
            try:
                tutor.update_skill_profile(user.id)
            except Exception:
                import traceback as tb
                tb.print_exc()
        threading.Thread(target=_bg_update_skills, daemon=True).start()

        return {
            "status": "recorded",
            "score": session.score,
            "total": session.total,
            "session_id": session.id,
        }
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# QUIZ SESSIONS
# ═══════════════════════════════════════════════════════════════════════════

class StartSessionRequest(BaseModel):
    paper_id: str
    paper_title: str = ""
    batch_number: int = 1


class SubmitAnswerRequest(BaseModel):
    question_id: str
    question_text: str = ""
    user_answer: str
    correct_answer: str = ""
    is_correct: bool
    time_spent_ms: int = 0


class CompleteSessionRequest(BaseModel):
    score: int
    total: int


@app.post("/api/sessions")
def start_session(req: StartSessionRequest, user: User = Depends(get_current_user)):
    db = SessionLocal()
    try:
        session = QuizSession(
            user_id=user.id,
            paper_id=req.paper_id,
            paper_title=req.paper_title,
            batch_number=req.batch_number,
        )
        db.add(session)
        db.commit()
        db.refresh(session)
        return {"session_id": session.id}
    finally:
        db.close()


@app.post("/api/sessions/{session_id}/answer")
def submit_answer(session_id: int, req: SubmitAnswerRequest, user: User = Depends(get_current_user)):
    db = SessionLocal()
    try:
        session = db.query(QuizSession).filter(QuizSession.id == session_id, QuizSession.user_id == user.id).first()
        if not session:
            raise HTTPException(404, "Session not found")

        answer = SessionAnswer(
            session_id=session_id,
            question_id=req.question_id,
            question_text=req.question_text,
            user_answer=req.user_answer,
            correct_answer=req.correct_answer,
            is_correct=req.is_correct,
            time_spent_ms=req.time_spent_ms,
        )
        db.add(answer)
        db.commit()
        return {"status": "ok"}
    finally:
        db.close()


@app.post("/api/sessions/{session_id}/complete")
def complete_session(session_id: int, req: CompleteSessionRequest, user: User = Depends(get_current_user)):
    db = SessionLocal()
    try:
        session = db.query(QuizSession).filter(QuizSession.id == session_id, QuizSession.user_id == user.id).first()
        if not session:
            raise HTTPException(404, "Session not found")

        session.score = req.score
        session.total = req.total
        session.completed = True
        session.completed_at = datetime.now(timezone.utc)
        db.commit()

        # Update skill profile in background (non-blocking)
        try:
            tutor.update_skill_profile(user.id)
        except Exception as e:
            print(f"[Skill Update] Warning: {e}")

        return {"status": "completed", "score": req.score, "total": req.total}
    finally:
        db.close()


@app.get("/api/sessions/history")
def get_session_history(user: User = Depends(get_current_user)):
    db = SessionLocal()
    try:
        sessions = (
            db.query(QuizSession)
            .filter(QuizSession.user_id == user.id)
            .order_by(QuizSession.started_at.desc())
            .limit(50)
            .all()
        )
        return [
            {
                "id": s.id,
                "paper_id": s.paper_id,
                "paper_title": s.paper_title,
                "score": s.score,
                "total": s.total,
                "batch": s.batch_number,
                "completed": s.completed,
                "started_at": s.started_at.isoformat() if s.started_at else None,
                "completed_at": s.completed_at.isoformat() if s.completed_at else None,
            }
            for s in sessions
        ]
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# ANSWER EVALUATION (LLM-powered grading for open-ended questions)
# ═══════════════════════════════════════════════════════════════════════════

class MarkPointInput(BaseModel):
    point: str
    marks: int = 1

class EvaluateAnswerRequest(BaseModel):
    question_text: str
    student_answer: str
    model_answer: Optional[str] = None
    key_terms: Optional[list[str]] = None
    mark_scheme: Optional[list[MarkPointInput]] = None
    total_marks: Optional[int] = None


GRADING_PROMPT_WITH_SCHEME = """You are a university exam marker. Grade the student's answer against the mark scheme.

QUESTION: {question_text}

MARK SCHEME:
{mark_scheme_text}

TOTAL MARKS: {total_marks}

MODEL ANSWER (for reference):
{model_answer}

STUDENT ANSWER:
{student_answer}

INSTRUCTIONS:
- Award marks for each marking point the student demonstrates, even if worded differently.
- A concise correct answer deserves full marks — do not penalise brevity.
- Focus on MEANING and conceptual accuracy, not exact wording.
- Be fair but rigorous — only award marks for points that are clearly addressed.

Respond with ONLY valid JSON (no markdown fences):
{{
  "marks_awarded": <int>,
  "total_marks": {total_marks},
  "points_hit": ["<marking point text that was addressed>", ...],
  "points_missed": ["<marking point text that was NOT addressed>", ...],
  "feedback": "<1-2 sentence constructive feedback>"
}}"""

GRADING_PROMPT_HOLISTIC = """You are a university exam marker. Grade the student's answer against the model answer.

QUESTION: {question_text}

MODEL ANSWER:
{model_answer}

KEY TERMS EXPECTED: {key_terms}

STUDENT ANSWER:
{student_answer}

INSTRUCTIONS:
- Award a score from 0 to 100 based on how well the student's answer captures the key concepts.
- A concise correct answer deserves a high score — do not penalise brevity.
- Focus on MEANING and conceptual accuracy, not exact wording or length.
- If the student uses different words to express the same concept, give credit.
- If key terms or their equivalents are missing, note them.

Respond with ONLY valid JSON (no markdown fences):
{{
  "score": <0-100>,
  "matched_terms": ["<terms the student addressed>", ...],
  "missing_terms": ["<important terms/concepts NOT addressed>", ...],
  "feedback": "<1-2 sentence constructive feedback>"
}}"""


@app.post("/api/evaluate-answer")
def evaluate_answer(req: EvaluateAnswerRequest):
    """Use GPT-4o-mini to evaluate an open-ended answer."""
    from openai import OpenAI

    api_key = os.getenv("OPENAI_API_KEY", "")
    if not api_key:
        raise HTTPException(500, "OpenAI API key not configured")

    client = OpenAI(api_key=api_key)

    has_scheme = req.mark_scheme and len(req.mark_scheme) > 0
    total_marks = req.total_marks or (sum(p.marks for p in req.mark_scheme) if has_scheme else 5)

    if has_scheme:
        scheme_lines = []
        for p in req.mark_scheme:
            scheme_lines.append(f"- [{p.marks} mark{'s' if p.marks != 1 else ''}] {p.point}")
        prompt = GRADING_PROMPT_WITH_SCHEME.format(
            question_text=req.question_text,
            mark_scheme_text="\n".join(scheme_lines),
            total_marks=total_marks,
            model_answer=req.model_answer or "(not provided)",
            student_answer=req.student_answer,
        )
    else:
        prompt = GRADING_PROMPT_HOLISTIC.format(
            question_text=req.question_text,
            model_answer=req.model_answer or "(not provided)",
            key_terms=", ".join(req.key_terms) if req.key_terms else "(none specified)",
            student_answer=req.student_answer,
        )

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=400,
            temperature=0.2,
        )
        raw = response.choices[0].message.content.strip()
        if raw.startswith("```"):
            lines = raw.split("\n")
            raw = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
        result = json.loads(raw)
    except Exception as e:
        raise HTTPException(500, f"Evaluation failed: {str(e)}")

    if has_scheme:
        marks = min(result.get("marks_awarded", 0), total_marks)
        return {
            "mode": "mark_scheme",
            "marks_awarded": marks,
            "total_marks": total_marks,
            "score": marks / total_marks if total_marks > 0 else 0,
            "is_correct": marks >= total_marks * 0.5,
            "points_hit": result.get("points_hit", []),
            "points_missed": result.get("points_missed", []),
            "feedback": result.get("feedback", ""),
        }
    else:
        score_pct = min(max(result.get("score", 0), 0), 100)
        return {
            "mode": "holistic",
            "score": score_pct / 100,
            "is_correct": score_pct >= 50,
            "matched_terms": result.get("matched_terms", []),
            "missing_terms": result.get("missing_terms", []),
            "feedback": result.get("feedback", ""),
        }


@app.get("/api/stats")
def get_user_stats(user: User = Depends(get_current_user)):
    """Aggregate stats for the user's dashboard, including streak."""
    from datetime import timedelta

    db = SessionLocal()
    try:
        sessions = db.query(QuizSession).filter(QuizSession.user_id == user.id, QuizSession.completed == True).all()
        total_sessions = len(sessions)
        total_questions = sum(s.total for s in sessions)
        total_correct = sum(s.score for s in sessions)
        avg_score = (total_correct / total_questions * 100) if total_questions > 0 else 0

        # Compute streak — consecutive days with at least one completed session
        today = datetime.now(timezone.utc).date()
        active_dates = set()
        for s in sessions:
            if s.completed_at:
                active_dates.add(s.completed_at.date())

        # Count streak backwards from today
        streak = 0
        check_date = today
        while check_date in active_dates:
            streak += 1
            check_date -= timedelta(days=1)

        # Build week activity (last 7 days, Mon-Sun aligned to current week)
        # Find the Monday of current week
        monday = today - timedelta(days=today.weekday())
        week_days = []
        day_labels = ['Mo', 'Tu', 'We', 'Th', 'Fr', 'Sa', 'Su']
        for i in range(7):
            d = monday + timedelta(days=i)
            if d > today:
                status = 'future'
            elif d in active_dates:
                status = 'active'
            else:
                status = 'missed'
            week_days.append({'label': day_labels[i], 'status': status})

        return {
            "total_sessions": total_sessions,
            "total_questions": total_questions,
            "total_correct": total_correct,
            "average_score": round(avg_score, 1),
            "streak": streak,
            "week": week_days,
        }
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# NOTEBOOKS
# ═══════════════════════════════════════════════════════════════════════════

@app.get("/api/notebooks")
def list_notebooks(user: User = Depends(get_current_user)):
    """List user's saved notebooks + all premade notebooks."""
    db = SessionLocal()
    try:
        notebooks = (
            db.query(SavedNotebook)
            .filter(
                (SavedNotebook.user_id == user.id) | (SavedNotebook.is_premade == True),
                SavedNotebook.deleted_at == None,
            )
            .order_by(SavedNotebook.created_at.desc())
            .all()
        )
        results = []
        for nb in notebooks:
            data = json.loads(nb.notebook_json)
            data["_saved_id"] = nb.id
            data["_is_premade"] = nb.is_premade
            data["_folder"] = nb.folder or ""
            results.append(data)
        return results
    finally:
        db.close()


@app.get("/api/notebooks/folders")
def list_folders(user: User = Depends(get_current_user)):
    """Return persisted folder names for this user."""
    db = SessionLocal()
    try:
        rows = db.query(StudyFolder).filter(StudyFolder.user_id == user.id).order_by(StudyFolder.created_at).all()
        return [r.name for r in rows]
    finally:
        db.close()


class MoveNotebookRequest(BaseModel):
    folder: str = ""


@app.put("/api/notebooks/{saved_id}/move")
def move_notebook(saved_id: int, req: MoveNotebookRequest, user: User = Depends(get_current_user)):
    """Move a notebook into a folder (or root if folder is empty)."""
    db = SessionLocal()
    try:
        nb = db.query(SavedNotebook).filter(
            SavedNotebook.id == saved_id, SavedNotebook.user_id == user.id
        ).first()
        if not nb:
            raise HTTPException(404, "Notebook not found")
        old_folder = nb.folder
        new_folder = req.folder.strip()[:100]
        nb.folder = new_folder
        db.commit()

        if old_folder and old_folder != new_folder:
            try:
                rag.delete_notebook_embeddings(user.id, old_folder, nb.notebook_id)
            except Exception:
                pass
        if new_folder:
            try:
                data = json.loads(nb.notebook_json)
                rag.embed_notebook(user.id, new_folder, nb.notebook_id, data)
            except Exception:
                import traceback
                traceback.print_exc()

        return {"status": "moved", "folder": nb.folder}
    finally:
        db.close()


class CreateFolderRequest(BaseModel):
    name: str


@app.post("/api/notebooks/folders")
def create_folder(req: CreateFolderRequest, user: User = Depends(get_current_user)):
    """Create and persist a new folder."""
    name = req.name.strip()[:100]
    if not name:
        raise HTTPException(400, "Folder name cannot be empty")
    db = SessionLocal()
    try:
        existing = db.query(StudyFolder).filter(
            StudyFolder.user_id == user.id, StudyFolder.name == name
        ).first()
        if existing:
            return {"folder": name}
        folder = StudyFolder(user_id=user.id, name=name)
        db.add(folder)
        db.commit()
        return {"folder": name}
    finally:
        db.close()


@app.delete("/api/notebooks/folders/{folder_name}")
def delete_folder(folder_name: str, user: User = Depends(get_current_user)):
    """Delete a folder and move its notebooks back to root."""
    db = SessionLocal()
    try:
        notebooks = db.query(SavedNotebook).filter(
            SavedNotebook.user_id == user.id, SavedNotebook.folder == folder_name
        ).all()
        for nb in notebooks:
            nb.folder = ""
        folder = db.query(StudyFolder).filter(
            StudyFolder.user_id == user.id, StudyFolder.name == folder_name
        ).first()
        if folder:
            db.delete(folder)
        db.commit()
        return {"status": "deleted", "moved_count": len(notebooks)}
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# FOLDER RAG
# ═══════════════════════════════════════════════════════════════════════════

@app.post("/api/folders/{folder_name}/embed")
def embed_folder(folder_name: str, user: User = Depends(get_current_user)):
    """Embed all notebooks in a folder into ChromaDB."""
    try:
        embed_uid = _curated_uid(folder_name) if _curated_uid(folder_name) is not None else user.id
        result = rag.embed_all_in_folder(embed_uid, folder_name)
        return result
    except Exception:
        import traceback
        traceback.print_exc()
        return {"notebooks_embedded": 0, "total_chunks": 0, "error": "Embedding failed"}

@app.post("/api/folders/{folder_name}/upload")
async def upload_folder_source(
    folder_name: str,
    file: UploadFile = File(...),
    authorization: Optional[str] = Header(None),
):
    """Upload a raw document to a folder — extract text, embed, no notebook generation."""
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(401, "Not authenticated")
    token = authorization.split(" ", 1)[1]
    payload = decode_access_token(token)
    if not payload:
        raise HTTPException(401, "Invalid token")
    db = SessionLocal()
    user = db.query(User).filter(User.id == int(payload["sub"])).first()
    db.close()
    if not user:
        raise HTTPException(401, "User not found")

    if not file.filename:
        raise HTTPException(400, "No filename provided")

    ext = Path(file.filename).suffix.lower()
    allowed = {".pdf", ".png", ".jpg", ".jpeg", ".pptx", ".tiff", ".bmp", ".webp"}
    if ext not in allowed:
        raise HTTPException(400, f"Unsupported file type: {ext}")

    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            content = await file.read()
            tmp.write(content)
            tmp_path = Path(tmp.name)

        raw_text = ""
        page_count = 0
        source_type = ext.lstrip(".")

        import asyncio

        def _extract_text_only():
            """Lightweight text extraction — no image loading to save memory."""
            if ext == ".pptx":
                from pptx import Presentation as PptxPres
                prs = PptxPres(str(tmp_path))
                pages = []
                for i, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                    pages.append({"slide_number": i + 1, "text": "\n".join(texts), "images": []})
                return pages
            elif ext == ".pdf":
                import pdfplumber
                pages = []
                with pdfplumber.open(str(tmp_path)) as pdf:
                    for i, pg in enumerate(pdf.pages):
                        txt = pg.extract_text() or ""
                        pages.append({"slide_number": i + 1, "text": txt, "images": []})
                return pages if pages else None
            return None

        pages = await asyncio.to_thread(_extract_text_only)

        if ext not in (".pdf", ".pptx"):
            return JSONResponse(status_code=400, content={"detail": "Image files must be uploaded via the full notebook pipeline"},
                                headers=_cors_headers)
        if pages:
            page_count = len(pages)
            raw_text = "\n\n".join(p.get("text", "") for p in pages if p.get("text"))
        else:
            return JSONResponse(status_code=400, content={"detail": "Could not extract text from this file"},
                                headers=_cors_headers)

        if not raw_text.strip():
            return JSONResponse(status_code=400, content={"detail": "No text could be extracted from this file"},
                                headers=_cors_headers)

        source_id = f"src_{uuid.uuid4().hex[:10]}"
        title = Path(file.filename).stem.replace("_", " ").replace("-", " ")

        stored_path = FOLDER_UPLOADS_DIR / f"{source_id}{ext}"
        shutil.copy2(str(tmp_path), str(stored_path))

        db = SessionLocal()
        try:
            fs = FolderSource(
                user_id=user.id,
                folder_name=folder_name,
                source_id=source_id,
                title=title,
                filename=file.filename,
                source_type=source_type,
                page_count=page_count,
                raw_text=raw_text,
                file_path=str(stored_path),
            )
            db.add(fs)
            db.commit()
        finally:
            db.close()

        def _bg_embed():
            try:
                rag.embed_raw_source(user.id, folder_name, source_id, title, raw_text)
                print(f"[BG] Embedded source {source_id} for folder {folder_name}")
            except Exception:
                import traceback
                traceback.print_exc()

        def _bg_images():
            from image_extractor import extract_and_store_images
            extract_and_store_images(stored_path, ext, source_id, user.id, folder_name, SOURCE_IMAGES_DIR)

        threading.Thread(target=_bg_embed, daemon=True).start()
        threading.Thread(target=_bg_images, daemon=True).start()

        return {
            "source_id": source_id,
            "title": title,
            "page_count": page_count,
            "filename": file.filename,
        }
    except HTTPException:
        raise
    except Exception:
        import traceback
        traceback.print_exc()
        return JSONResponse(status_code=500, content={"detail": "Upload failed — server error"},
                            headers=_cors_headers)
    finally:
        if tmp_path:
            tmp_path.unlink(missing_ok=True)


@app.delete("/api/folders/{folder_name}/sources/{source_id}")
def delete_folder_source(folder_name: str, source_id: str, user: User = Depends(get_current_user)):
    """Delete a raw source from a folder."""
    owner_id = _curated_uid(folder_name) if _curated_uid(folder_name) is not None else user.id
    db = SessionLocal()
    file_path = None
    try:
        fs = db.query(FolderSource).filter(
            FolderSource.user_id == owner_id,
            FolderSource.folder_name == folder_name,
            FolderSource.source_id == source_id,
        ).first()
        if not fs:
            raise HTTPException(404, "Source not found")
        file_path = fs.file_path
        db.delete(fs)
        db.commit()
    finally:
        db.close()

    if file_path:
        Path(file_path).unlink(missing_ok=True)

    try:
        rag.delete_notebook_embeddings(owner_id, folder_name, source_id)
    except Exception:
        pass

    return {"status": "deleted", "source_id": source_id}


@app.get("/api/folders/{folder_name}/sources/{source_id}/file")
def get_source_file(folder_name: str, source_id: str, user: User = Depends(get_current_user)):
    """Serve the original uploaded file for a folder source."""
    owner_id = _curated_uid(folder_name) if _curated_uid(folder_name) is not None else user.id
    db = SessionLocal()
    try:
        fs = db.query(FolderSource).filter(
            FolderSource.user_id == owner_id,
            FolderSource.folder_name == folder_name,
            FolderSource.source_id == source_id,
        ).first()
        if not fs:
            raise HTTPException(404, "Source not found")
        if not fs.file_path or not Path(fs.file_path).exists():
            raise HTTPException(404, "File not available")
        file_path = Path(fs.file_path)
        filename = fs.filename
    finally:
        db.close()

    media_types = {
        ".pdf": "application/pdf",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    media_type = media_types.get(file_path.suffix.lower(), "application/octet-stream")

    from fastapi.responses import FileResponse
    return FileResponse(
        path=str(file_path),
        media_type=media_type,
        filename=filename,
        headers={"Content-Disposition": f'inline; filename="{filename}"'},
    )


def _resolve_image_path(stored_path: str) -> Path | None:
    """Find the actual image file, trying persistent disk if stored path is stale."""
    fp = Path(stored_path)
    if fp.exists():
        return fp
    _OLD_PREFIX = "/opt/render/project/src/folder_uploads"
    _NEW_PREFIX = "/data/folder_uploads"
    if stored_path.startswith(_OLD_PREFIX):
        alt = Path(_NEW_PREFIX + stored_path[len(_OLD_PREFIX):])
        if alt.exists():
            return alt
    elif stored_path.startswith(_NEW_PREFIX):
        alt = Path(_OLD_PREFIX + stored_path[len(_NEW_PREFIX):])
        if alt.exists():
            return alt
    return None


@app.get("/api/source-images/{image_id}")
def serve_source_image(image_id: int):
    """Serve an extracted source image by its DB id."""
    db = SessionLocal()
    try:
        si = db.query(SourceImage).filter(SourceImage.id == image_id).first()
        if not si:
            raise HTTPException(404, "Image not found in DB")
        resolved = _resolve_image_path(si.image_path)
        if not resolved:
            print(f"[images] File missing: {si.image_path}  (id={image_id}, source={si.source_id})")
            raise HTTPException(404, "Image file missing")
        if str(resolved) != si.image_path:
            si.image_path = str(resolved)
            db.commit()
        from fastapi.responses import FileResponse
        return FileResponse(
            path=str(resolved),
            media_type="image/png",
            headers={"Cache-Control": "public, max-age=86400"},
        )
    finally:
        db.close()


@app.get("/api/debug/images")
def debug_images():
    """Debug endpoint: show image status and curated source file_paths."""
    db = SessionLocal()
    try:
        all_imgs = db.query(SourceImage).all()
        by_folder: dict[str, dict] = {}
        for si in all_imgs:
            f = si.folder_name
            if f not in by_folder:
                by_folder[f] = {"total": 0, "exists": 0, "missing": 0, "sample_ids": [], "sample_paths": []}
            by_folder[f]["total"] += 1
            if Path(si.image_path).exists():
                by_folder[f]["exists"] += 1
            else:
                by_folder[f]["missing"] += 1
            if len(by_folder[f]["sample_ids"]) < 3:
                by_folder[f]["sample_ids"].append(si.id)
                by_folder[f]["sample_paths"].append(si.image_path)

        curated_sources = []
        from curated_config import CURATED_FOLDER_NAMES, CURATED_USER_ID
        for fn in CURATED_FOLDER_NAMES:
            sources = db.query(FolderSource).filter(
                FolderSource.user_id == CURATED_USER_ID,
                FolderSource.folder_name == fn,
            ).all()
            for s in sources[:3]:
                fp = Path(s.file_path) if s.file_path else None
                curated_sources.append({
                    "folder": fn,
                    "source_id": s.source_id,
                    "title": s.title[:50],
                    "file_path": s.file_path,
                    "file_exists": fp.exists() if fp else False,
                })

        data_dir = Path("/data")
        return {
            "images_by_folder": by_folder,
            "curated_sources_sample": curated_sources,
            "persistent_disk_exists": data_dir.is_dir(),
            "data_contents": [str(p) for p in data_dir.iterdir()] if data_dir.is_dir() else [],
        }
    finally:
        db.close()


@app.get("/api/folders/{folder_name}/images")
def list_folder_images(folder_name: str, user: User = Depends(get_current_user)):
    """List all extracted images for a folder."""
    owner_id = _curated_uid(folder_name) if _curated_uid(folder_name) is not None else user.id
    db = SessionLocal()
    try:
        images = db.query(SourceImage).filter(
            SourceImage.user_id == owner_id,
            SourceImage.folder_name == folder_name,
        ).all()
        return [
            {
                "id": si.id,
                "source_id": si.source_id,
                "page_number": si.page_number,
                "context_text": si.context_text[:200],
                "width": si.width,
                "height": si.height,
            }
            for si in images
        ]
    finally:
        db.close()


@app.post("/api/folders/{folder_name}/sources/{source_id}/generate-notebook")
async def generate_notebook_from_source(
    folder_name: str,
    source_id: str,
    authorization: Optional[str] = Header(None),
):
    """Run the full notebook pipeline on a stored folder source, returning SSE stream."""
    import asyncio
    import queue
    import threading

    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(401, "Not authenticated")
    token = authorization.split(" ", 1)[1]
    payload = decode_access_token(token)
    if not payload:
        raise HTTPException(401, "Invalid token")
    db = SessionLocal()
    user = db.query(User).filter(User.id == int(payload["sub"])).first()
    fs = db.query(FolderSource).filter(
        FolderSource.source_id == source_id,
        FolderSource.folder_name == folder_name,
    ).first() if user else None
    db.close()
    if not user:
        raise HTTPException(401, "User not found")
    if not fs:
        raise HTTPException(404, "Source not found")
    if not fs.file_path or not Path(fs.file_path).exists():
        raise HTTPException(404, "Original file not available")

    if user:
        usage = _get_user_usage(user.id)
        if usage["notebooks_remaining"] <= 0:
            raise HTTPException(429, "Notebook limit reached.")

    src_path = Path(fs.file_path)
    progress_q: queue.Queue = queue.Queue()

    def on_progress(stage, current, total):
        msg = _STAGE_MESSAGES.get(stage, stage)
        try:
            msg = msg.format(current=current, total=total)
        except (KeyError, IndexError):
            pass
        progress_q.put({"stage": stage, "message": msg, "current": current, "total": total})

    def run_pipeline():
        try:
            from pipeline import run_notebook_pipeline
            from curated_config import get_course_for_folder
            course = get_course_for_folder(folder_name)
            paper_paths = _get_paper_paths(course)
            result = run_notebook_pipeline(
                src_path,
                output_path=None,
                paper_paths=paper_paths if paper_paths else None,
                provider="openai",
                extra_instructions="",
                detail="low",
                on_progress=on_progress,
            )
            nb_id = result.get("id", f"nb_{uuid.uuid4().hex[:8]}")
            save_path = GENERATED_DIR / f"{nb_id}.json"
            with open(save_path, "w", encoding="utf-8") as f:
                json.dump(result, f, indent=2, ensure_ascii=False)

            db2 = SessionLocal()
            try:
                title = result.get("title", "Untitled")
                saved = SavedNotebook(
                    user_id=user.id,
                    notebook_id=nb_id,
                    title=title,
                    course=result.get("course", ""),
                    notebook_json=json.dumps(result),
                    is_premade=False,
                    folder=folder_name,
                )
                db2.add(saved)
                db2.commit()
                db2.refresh(saved)
                result["_saved_id"] = saved.id
                result["_folder"] = folder_name

                threading.Thread(
                    target=_bg_post_process,
                    args=(user.id, folder_name, nb_id, result),
                    daemon=True,
                ).start()
            finally:
                db2.close()

            progress_q.put({"stage": "done", "notebook": result})
        except Exception as exc:
            progress_q.put({"stage": "error", "message": str(exc)})

    threading.Thread(target=run_pipeline, daemon=True).start()

    async def event_stream():
        idle_ticks = 0
        while True:
            try:
                item = progress_q.get_nowait()
                idle_ticks = 0
            except queue.Empty:
                await asyncio.sleep(0.5)
                idle_ticks += 1
                if idle_ticks % 20 == 0:
                    yield ": keepalive\n\n"
                continue
            yield f"data: {json.dumps(item)}\n\n"
            if item.get("stage") in ("done", "error"):
                break

    return StreamingResponse(event_stream(), media_type="text/event-stream")


@app.get("/api/folders/{folder_name}/sources")
def folder_sources(folder_name: str, user: User = Depends(get_current_user)):
    """List all sources (notebooks + raw documents) with metadata."""
    src_uid = _curated_uid(folder_name) if _curated_uid(folder_name) is not None else user.id
    try:
        chroma_sources = rag.get_folder_sources(src_uid, folder_name)
    except Exception:
        import traceback
        traceback.print_exc()
        chroma_sources = []

    embedded_ids = {s["notebook_id"] for s in chroma_sources}

    db = SessionLocal()
    try:
        notebooks = (
            db.query(SavedNotebook)
            .filter(
                SavedNotebook.user_id == src_uid,
                SavedNotebook.folder == folder_name,
                SavedNotebook.deleted_at == None,
            )
            .all()
        )
        raw_sources = (
            db.query(FolderSource)
            .filter(
                FolderSource.user_id == src_uid,
                FolderSource.folder_name == folder_name,
            )
            .all()
        )

        all_sources = []
        for nb in notebooks:
            data = json.loads(nb.notebook_json)
            all_sources.append({
                "notebook_id": nb.notebook_id,
                "saved_id": nb.id,
                "title": data.get("title", nb.title),
                "course": data.get("course", nb.course),
                "section_count": len(data.get("sections") or []),
                "embedded": nb.notebook_id in embedded_ids,
                "type": "notebook",
            })
        for src in raw_sources:
            all_sources.append({
                "notebook_id": src.source_id,
                "source_id": src.source_id,
                "title": src.title,
                "filename": src.filename,
                "source_type": src.source_type,
                "page_count": src.page_count,
                "embedded": src.source_id in embedded_ids,
                "type": "document",
            })

        return {"sources": all_sources, "embedding_stats": chroma_sources}
    finally:
        db.close()

@app.post("/api/folders/{folder_name}/study-plan")
def folder_study_plan(folder_name: str, user: User = Depends(get_current_user)):
    """Generate a study plan from all sources in a folder."""
    plan = rag.generate_study_plan(user.id, folder_name, user.name)
    return {"plan": plan}

@app.get("/api/folders/{folder_name}/notebooks")
def folder_notebooks(folder_name: str, user: User = Depends(get_current_user)):
    """List all notebooks in a folder."""
    db = SessionLocal()
    try:
        notebooks = (
            db.query(SavedNotebook)
            .filter(
                SavedNotebook.user_id == user.id,
                SavedNotebook.folder == folder_name,
                SavedNotebook.deleted_at == None,
            )
            .all()
        )
        return [
            {
                "id": nb.notebook_id,
                "_saved_id": nb.id,
                "title": nb.title,
                "course": nb.course,
                **json.loads(nb.notebook_json),
            }
            for nb in notebooks
        ]
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# LESSON / COURSE OUTLINE
# ═══════════════════════════════════════════════════════════════════════════

@app.post("/api/folders/{folder_name}/outline")
def generate_outline(folder_name: str, user: User = Depends(get_current_user)):
    """Generate or regenerate a course outline from folder sources."""
    src_uid = _curated_uid(folder_name) if _curated_uid(folder_name) is not None else user.id
    structure = get_lesson_structure(folder_name)
    result = lesson.generate_outline(user.id, folder_name, source_user_id=src_uid, structure=structure)
    if "error" in result:
        raise HTTPException(400, result["error"])
    return result


@app.get("/api/folders/{folder_name}/lesson")
def get_lesson_state(folder_name: str, user: User = Depends(get_current_user)):
    """Get current lesson state — outline, progress, current section."""
    return lesson.get_lesson_state(user.id, folder_name)


@app.post("/api/folders/{folder_name}/lesson/advance")
def advance_lesson(folder_name: str, user: User = Depends(get_current_user)):
    """Advance to the next lesson section."""
    result = lesson.advance_section(user.id, folder_name)
    if "error" in result:
        raise HTTPException(400, result["error"])
    return result


@app.post("/api/folders/{folder_name}/lesson/reset")
def reset_lesson(folder_name: str, user: User = Depends(get_current_user)):
    """Reset lesson progress to start over."""
    result = lesson.reset_lesson(user.id, folder_name)
    if "error" in result:
        raise HTTPException(400, result["error"])
    return result


@app.get("/api/folders/{folder_name}/lesson-notes")
def get_lesson_notes(folder_name: str, user: User = Depends(get_current_user)):
    db = SessionLocal()
    try:
        row = db.query(LessonNotes).filter(
            LessonNotes.user_id == user.id,
            LessonNotes.folder_name == folder_name,
        ).first()
        return {"content_html": row.content_html if row else ""}
    finally:
        db.close()


@app.put("/api/folders/{folder_name}/lesson-notes")
def save_lesson_notes(folder_name: str, body: dict, user: User = Depends(get_current_user)):
    html = body.get("content_html", "")
    db = SessionLocal()
    try:
        row = db.query(LessonNotes).filter(
            LessonNotes.user_id == user.id,
            LessonNotes.folder_name == folder_name,
        ).first()
        if row:
            row.content_html = html
            row.updated_at = datetime.now(timezone.utc)
        else:
            row = LessonNotes(
                user_id=user.id,
                folder_name=folder_name,
                content_html=html,
            )
            db.add(row)
        db.commit()
        return {"ok": True}
    finally:
        db.close()


@app.get("/api/folders/{folder_name}/all-feedback")
def get_all_feedback(folder_name: str, user: User = Depends(get_current_user)):
    return {"sections": []}


@app.get("/api/folders/{folder_name}/section-chat/{section_index}")
def get_section_chat(folder_name: str, section_index: int, user: User = Depends(get_current_user)):
    db = SessionLocal()
    try:
        rows = (
            db.query(ChatMessage)
            .filter(
                ChatMessage.user_id == user.id,
                ChatMessage.context_type == "lesson",
                ChatMessage.context_id == folder_name,
                ChatMessage.section_index == section_index,
            )
            .order_by(ChatMessage.created_at.asc())
            .all()
        )
        messages = [{"role": r.role, "content": r.content} for r in rows]
        return {"messages": messages}
    finally:
        db.close()


@app.post("/api/notebooks/save")
def save_notebook(notebook: dict, user: User = Depends(get_current_user)):
    """Save a generated notebook to the user's account."""
    db = SessionLocal()
    try:
        nb_id = notebook.get("id", f"nb_{uuid.uuid4().hex[:8]}")
        saved = SavedNotebook(
            user_id=user.id,
            notebook_id=nb_id,
            title=notebook.get("title", "Untitled"),
            course=notebook.get("course", ""),
            notebook_json=json.dumps(notebook),
            is_premade=False,
            folder=notebook.get("folder", ""),
        )
        db.add(saved)
        db.commit()
        db.refresh(saved)

        folder = notebook.get("folder", "")
        threading.Thread(
            target=_bg_post_process,
            args=(user.id, folder, nb_id, notebook),
            daemon=True,
        ).start()

        return {"status": "saved", "id": saved.id}
    finally:
        db.close()


@app.delete("/api/notebooks/{notebook_id}")
def delete_notebook(notebook_id: int, user: User = Depends(get_current_user)):
    print(f"  [delete] user={user.id} notebook_id={notebook_id}")
    db = SessionLocal()
    try:
        nb = db.query(SavedNotebook).filter(SavedNotebook.id == notebook_id, SavedNotebook.user_id == user.id).first()
        if not nb:
            print(f"  [delete] NOT FOUND — id={notebook_id} user={user.id}")
            raise HTTPException(404, "Notebook not found")

        slug = nb.notebook_id
        title = nb.title[:40]
        db.delete(nb)
        db.commit()

        print(f"  [delete] HARD-DELETED — id={notebook_id} slug={slug} title={title}")
        return {"status": "deleted", "notebook_id": slug}
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# PEDRO CHAT (AI TUTOR)
# ═══════════════════════════════════════════════════════════════════════════

class ChatSendRequest(BaseModel):
    message: str
    conversation_id: Optional[str] = None
    context_type: str = "global"  # "notebook", "global", "session"
    context_id: Optional[str] = None
    notebook_ids: Optional[list[str]] = None
    section_index: Optional[int] = None


@app.post("/api/chat/send")
def chat_send(req: ChatSendRequest, user: User = Depends(get_current_user)):
    """Send a message to Pedro and get a Socratic response."""
    if not req.message.strip():
        raise HTTPException(400, "Message cannot be empty")
    if req.context_type not in ("notebook", "global", "session", "folder", "lesson"):
        raise HTTPException(400, "Invalid context_type")

    usage = _get_user_usage(user.id)
    if usage["chat_messages_remaining"] <= 0:
        raise HTTPException(
            429,
            "You've reached your weekly message limit. "
            "Your limit resets in a few days — thanks for testing Coast!",
        )

    try:
        result = tutor.send_message(
            user_id=user.id,
            message=req.message.strip(),
            conversation_id=req.conversation_id,
            context_type=req.context_type,
            context_id=req.context_id,
            notebook_ids=req.notebook_ids,
        )
        result["usage"] = _get_user_usage(user.id)
        return result
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"Chat error: {str(e)}")


@app.post("/api/chat/stream")
def chat_stream(req: ChatSendRequest, user: User = Depends(get_current_user)):
    """Streaming version of chat/send — returns SSE with token chunks."""
    if not req.message.strip():
        raise HTTPException(400, "Message cannot be empty")
    if req.context_type not in ("notebook", "global", "session", "folder", "lesson"):
        raise HTTPException(400, "Invalid context_type")

    usage = _get_user_usage(user.id)
    if usage["chat_messages_remaining"] <= 0:
        raise HTTPException(429, "Weekly message limit reached.")

    def event_stream():
        try:
            for token, meta in tutor.send_message_stream(
                user_id=user.id,
                message=req.message.strip(),
                conversation_id=req.conversation_id,
                context_type=req.context_type,
                context_id=req.context_id,
                notebook_ids=req.notebook_ids,
                section_index=req.section_index,
            ):
                if token is not None:
                    yield f"data: {json.dumps({'token': token})}\n\n"
                if meta is not None:
                    meta["usage"] = _get_user_usage(user.id)
                    yield f"data: {json.dumps({'done': True, **meta})}\n\n"
        except Exception as e:
            import traceback
            traceback.print_exc()
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")


@app.get("/api/chat/history")
def chat_history(conversation_id: str, user: User = Depends(get_current_user)):
    """Get all messages for a specific conversation."""
    messages = tutor.get_chat_history(conversation_id, user.id)
    return messages


@app.get("/api/chat/conversations")
def chat_conversations(user: User = Depends(get_current_user)):
    """List all of a user's conversations with Pedro."""
    return tutor.get_conversations(user.id)


class AddNoteRequest(BaseModel):
    pedro_message: str
    notebook_id: Optional[str] = None


class ExerciseRequest(BaseModel):
    section_title: str
    section_content: str
    action: str = "generate"  # "generate" or "evaluate"
    question: str = ""
    answer: str = ""


@app.post("/api/exercise")
def exercise(req: ExerciseRequest, user: User = Depends(get_current_user)):
    """Generate a practice question or evaluate a student's answer."""
    try:
        result = tutor.handle_exercise(
            user_id=user.id,
            section_title=req.section_title,
            section_content=req.section_content,
            action=req.action,
            question=req.question,
            student_answer=req.answer,
        )
        return result
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"Exercise error: {str(e)}")


@app.post("/api/chat/add-note")
def chat_add_note(req: AddNoteRequest, user: User = Depends(get_current_user)):
    """Condense a Pedro message into a study note for the notebook."""
    if not req.pedro_message.strip():
        raise HTTPException(400, "Message cannot be empty")

    try:
        note_html = tutor.generate_note_for_notebook(req.pedro_message.strip())
        return {"note_html": note_html}
    except Exception as e:
        raise HTTPException(500, f"Note generation error: {str(e)}")


@app.get("/api/skill-profile")
def skill_profile(user: User = Depends(get_current_user)):
    """Get the user's topic skill profile."""
    return tutor.get_skill_profile(user.id)


@app.get("/api/tutor-memo")
def tutor_memo_endpoint(user: User = Depends(get_current_user)):
    """Get Pedro's memo about the user (for transparency / debugging)."""
    return tutor.get_tutor_memo(user.id)


# ═══════════════════════════════════════════════════════════════════════════
# SPACED REPETITION
# ═══════════════════════════════════════════════════════════════════════════

class ConceptExtractRequest(BaseModel):
    notebook_id: str

class ReviewSubmitRequest(BaseModel):
    card_id: int
    quality: int  # 0-5

@app.post("/api/concepts/extract")
def extract_concepts_endpoint(req: ConceptExtractRequest, user: User = Depends(get_current_user)):
    """Extract review concepts from a saved notebook."""
    db = SessionLocal()
    try:
        nb = db.query(SavedNotebook).filter(
            SavedNotebook.notebook_id == req.notebook_id,
            SavedNotebook.user_id == user.id,
        ).first()
        if not nb:
            raise HTTPException(status_code=404, detail="Notebook not found")
        notebook_json = json.loads(nb.notebook_json)
        count = spaced_rep.create_cards_for_notebook(user.id, req.notebook_id, notebook_json)
        return {"status": "ok", "cards_created": count}
    finally:
        db.close()

@app.get("/api/review/due")
def review_due(user: User = Depends(get_current_user)):
    """Get review cards due now."""
    cards = spaced_rep.get_due_cards(user.id)
    return {"cards": cards, "count": len(cards)}

@app.post("/api/review/submit")
def review_submit(req: ReviewSubmitRequest, user: User = Depends(get_current_user)):
    """Submit a review quality grade and update SM-2 schedule."""
    result = spaced_rep.submit_review(req.card_id, user.id, req.quality)
    if "error" in result:
        raise HTTPException(status_code=400, detail=result["error"])
    return result

@app.get("/api/review/stats")
def review_stats(user: User = Depends(get_current_user)):
    """Get spaced repetition stats for the dashboard."""
    return spaced_rep.get_review_stats(user.id)

@app.get("/api/dashboard/briefing")
def dashboard_briefing(user: User = Depends(get_current_user)):
    """Get Pedro's personalized daily briefing."""
    message = spaced_rep.generate_briefing(user.id, user.name)
    stats = spaced_rep.get_review_stats(user.id)
    return {"message": message, "review_stats": stats}


# ═══════════════════════════════════════════════════════════════════════════
# GENERATE NOTES (OCR PIPELINE)
# ═══════════════════════════════════════════════════════════════════════════

_STAGE_MESSAGES = {
    "extracting": "Extracting text and images...",
    "loaded":     "Found {total} pages",
    "chunking":   "Splitting into {total} chunks...",
    "analyzing":  "Analyzing chunk {current} of {total} with AI...",
    "merging":    "Merging sections into final guide...",
    "matching":   "Matching past paper questions...",
}


@app.post("/api/generate-notes")
async def generate_notes(
    file: UploadFile = File(...),
    instructions: str = Form(""),
    provider: str = Form("openai"),
    detail: str = Form("low"),
    folder: str = Form(""),
    authorization: Optional[str] = Header(None),
):
    """Upload a PDF/image and stream progress via SSE, final event has the notebook."""
    import asyncio
    import queue
    import threading

    user = None
    if authorization and authorization.startswith("Bearer "):
        token = authorization.split(" ", 1)[1]
        payload = decode_access_token(token)
        if payload:
            db = SessionLocal()
            user = db.query(User).filter(User.id == int(payload["sub"])).first()
            db.close()

    if user:
        usage = _get_user_usage(user.id)
        if usage["notebooks_remaining"] <= 0:
            raise HTTPException(
                429,
                f"You've reached your notebook limit ({RATE_LIMIT_NOTEBOOKS} notebooks). "
                "Delete an existing notebook to free up a slot.",
            )

    if not file.filename:
        raise HTTPException(400, "No filename provided")

    ext = Path(file.filename).suffix.lower()
    allowed = {".pdf", ".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp", ".gif", ".pptx"}
    if ext not in allowed:
        raise HTTPException(400, f"Unsupported file type: {ext}")

    if detail not in ("low", "high"):
        detail = "low"
    if provider not in ("openai", "anthropic", "kimi"):
        provider = "openai"

    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        content = await file.read()
        tmp.write(content)
        tmp_path = Path(tmp.name)

    progress_q: queue.Queue = queue.Queue()

    def on_progress(stage: str, current: int, total: int):
        msg = _STAGE_MESSAGES.get(stage, stage)
        try:
            msg = msg.format(current=current, total=total)
        except (KeyError, IndexError):
            pass
        progress_q.put({"stage": stage, "message": msg, "current": current, "total": total})

    def run_pipeline_thread():
        try:
            from pipeline import run_notebook_pipeline
            from curated_config import get_course_for_folder

            course = get_course_for_folder(folder) if folder else None
            paper_paths = _get_paper_paths(course)
            result = run_notebook_pipeline(
                tmp_path,
                output_path=None,
                paper_paths=paper_paths if paper_paths else None,
                provider=provider,
                extra_instructions=instructions,
                detail=detail,
                on_progress=on_progress,
            )

            nb_id = result.get("id", f"nb_{uuid.uuid4().hex[:8]}")
            save_path = GENERATED_DIR / f"{nb_id}.json"
            counter = 1
            while save_path.exists():
                save_path = GENERATED_DIR / f"{nb_id}_{counter}.json"
                counter += 1

            with open(save_path, "w", encoding="utf-8") as f:
                json.dump(result, f, indent=2, ensure_ascii=False)

            if user:
                db = SessionLocal()
                try:
                    from datetime import timedelta
                    title = result.get("title", "Untitled")
                    recent_cutoff = datetime.now(timezone.utc) - timedelta(seconds=60)
                    duplicate = (
                        db.query(SavedNotebook)
                        .filter(
                            SavedNotebook.user_id == user.id,
                            SavedNotebook.title == title,
                            SavedNotebook.created_at >= recent_cutoff,
                        )
                        .first()
                    )
                    if duplicate:
                        result["_saved_id"] = duplicate.id
                        result["_folder"] = duplicate.folder or ""
                    else:
                        target_folder = folder.strip()[:100] if folder else ""
                        saved = SavedNotebook(
                            user_id=user.id,
                            notebook_id=nb_id,
                            title=title,
                            course=result.get("course", ""),
                            notebook_json=json.dumps(result),
                            is_premade=False,
                            folder=target_folder,
                        )
                        db.add(saved)
                        db.commit()
                        db.refresh(saved)
                        result["_saved_id"] = saved.id
                        result["_folder"] = target_folder

                        threading.Thread(
                            target=_bg_post_process,
                            args=(user.id, target_folder, nb_id, result),
                            daemon=True,
                        ).start()
                finally:
                    db.close()

            progress_q.put({"stage": "done", "notebook": result})
        except Exception as exc:
            progress_q.put({"stage": "error", "message": str(exc)})
        finally:
            tmp_path.unlink(missing_ok=True)
            notebook_dir = tmp_path.parent / f"{tmp_path.stem}_notebook"
            if notebook_dir.exists():
                shutil.rmtree(notebook_dir, ignore_errors=True)

    async def event_stream():
        loop = asyncio.get_event_loop()
        threading.Thread(target=run_pipeline_thread, daemon=True).start()

        while True:
            try:
                msg = await loop.run_in_executor(None, lambda: progress_q.get(timeout=5))
            except queue.Empty:
                yield ": keepalive\n\n"
                continue

            yield f"data: {json.dumps(msg)}\n\n"

            if msg.get("stage") in ("done", "error"):
                break

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


# ═══════════════════════════════════════════════════════════════════════════
# ADMIN DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════


@app.get("/api/admin/overview")
def admin_overview(user: User = Depends(get_current_user)):
    """Return all users with stats, skill profiles, and tutor memos. Admin only."""
    if user.email != ADMIN_EMAIL:
        raise HTTPException(403, "Admin access only")

    from datetime import timedelta

    db = SessionLocal()
    try:
        users = db.query(User).order_by(User.created_at.desc()).all()
        result = []

        for u in users:
            # Sessions
            sessions = db.query(QuizSession).filter(
                QuizSession.user_id == u.id, QuizSession.completed == True
            ).all()
            total_q = sum(s.total for s in sessions)
            total_correct = sum(s.score for s in sessions)

            # Streak
            today = datetime.now(timezone.utc).date()
            active_dates = set()
            for s in sessions:
                if s.completed_at:
                    active_dates.add(s.completed_at.date())
            streak = 0
            check = today
            while check in active_dates:
                streak += 1
                check -= timedelta(days=1)

            # Skill profile
            sp = db.query(SkillProfile).filter(SkillProfile.user_id == u.id).first()
            skill = json.loads(sp.profile_json) if sp else {}

            # Tutor memo
            memo = db.query(TutorMemo).filter(TutorMemo.user_id == u.id).first()

            # Chat message count
            msg_count = db.query(ChatMessage).filter(ChatMessage.user_id == u.id).count()

            # Notebooks count
            nb_count = db.query(SavedNotebook).filter(
                SavedNotebook.user_id == u.id, SavedNotebook.is_premade == False
            ).count()

            # Folders / lessons / sections
            from database import CourseOutline
            folder_count = db.query(StudyFolder).filter(StudyFolder.user_id == u.id).count()
            user_outlines = db.query(CourseOutline).filter(CourseOutline.user_id == u.id).all()
            lessons_started = len(user_outlines)
            sections_done = sum(o.current_section for o in user_outlines)

            last_msg = (
                db.query(ChatMessage)
                .filter(ChatMessage.user_id == u.id)
                .order_by(ChatMessage.created_at.desc())
                .first()
            )
            last_active = (
                last_msg.created_at.isoformat()
                if last_msg and last_msg.created_at
                else (u.created_at.isoformat() if u.created_at else None)
            )

            result.append({
                "id": u.id,
                "name": u.name,
                "email": u.email,
                "course": u.course,
                "created_at": u.created_at.isoformat() if u.created_at else None,
                "sessions_completed": len(sessions),
                "total_questions": total_q,
                "total_correct": total_correct,
                "accuracy": round(total_correct / total_q * 100, 1) if total_q > 0 else 0,
                "streak": streak,
                "skill_profile": skill,
                "tutor_memo": memo.memo_text if memo else "",
                "memo_updated_at": memo.updated_at.isoformat() if memo and memo.updated_at else None,
                "chat_messages": msg_count,
                "notebooks_generated": nb_count,
                "folders_created": folder_count,
                "lessons_started": lessons_started,
                "sections_completed": sections_done,
                "last_active": last_active,
            })

        return {"users": result, "total_users": len(result)}
    finally:
        db.close()


@app.get("/api/admin/export-cohort")
def export_cohort(user: User = Depends(get_current_user)):
    """Full data export for the current cohort — quiz answers, chat logs, skill profiles."""
    if user.email != ADMIN_EMAIL:
        raise HTTPException(403, "Admin access only")

    db = SessionLocal()
    try:
        users = db.query(User).all()
        export = []

        for u in users:
            if u.email == ADMIN_EMAIL:
                continue

            sessions = db.query(QuizSession).filter(QuizSession.user_id == u.id).all()
            session_data = []
            for s in sessions:
                answers = db.query(SessionAnswer).filter(SessionAnswer.session_id == s.id).all()
                session_data.append({
                    "paper_id": s.paper_id,
                    "paper_title": s.paper_title,
                    "score": s.score,
                    "total": s.total,
                    "completed": s.completed,
                    "started_at": s.started_at.isoformat() if s.started_at else None,
                    "completed_at": s.completed_at.isoformat() if s.completed_at else None,
                    "answers": [
                        {
                            "question_id": a.question_id,
                            "question_text": a.question_text,
                            "user_answer": a.user_answer,
                            "correct_answer": a.correct_answer,
                            "is_correct": a.is_correct,
                            "time_spent_ms": a.time_spent_ms,
                        }
                        for a in answers
                    ],
                })

            messages = (
                db.query(ChatMessage)
                .filter(ChatMessage.user_id == u.id)
                .order_by(ChatMessage.created_at)
                .all()
            )
            chat_data = [
                {
                    "conversation_id": m.conversation_id,
                    "role": m.role,
                    "content": m.content,
                    "context_type": m.context_type,
                    "created_at": m.created_at.isoformat() if m.created_at else None,
                }
                for m in messages
            ]

            sp = db.query(SkillProfile).filter(SkillProfile.user_id == u.id).first()
            memo = db.query(TutorMemo).filter(TutorMemo.user_id == u.id).first()

            nb_count = (
                db.query(SavedNotebook)
                .filter(SavedNotebook.user_id == u.id, SavedNotebook.is_premade == False)
                .count()
            )

            export.append({
                "user_id": u.id,
                "name": u.name,
                "email": u.email,
                "course": u.course,
                "created_at": u.created_at.isoformat() if u.created_at else None,
                "skill_profile": json.loads(sp.profile_json) if sp else {},
                "tutor_memo": memo.memo_text if memo else "",
                "notebooks_generated": nb_count,
                "quiz_sessions": session_data,
                "chat_messages": chat_data,
            })

        return {
            "exported_at": datetime.now(timezone.utc).isoformat(),
            "total_students": len(export),
            "students": export,
        }
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# ACTIVITY TRACKING
# ═══════════════════════════════════════════════════════════════════════════


class ActivityRequest(BaseModel):
    feature: str
    duration_ms: int = 0
    action: str = "session"


@app.post("/api/activity")
def log_activity(body: ActivityRequest, user: User = Depends(get_current_user)):
    """Log a feature usage event (time spent on a feature session)."""
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    db = SessionLocal()
    try:
        ev = ActivityEvent(
            user_id=user.id,
            feature=body.feature,
            action=body.action,
            duration_ms=max(body.duration_ms, 0),
            event_date=today,
        )
        db.add(ev)
        db.commit()
        return {"ok": True}
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# ANALYTICS
# ═══════════════════════════════════════════════════════════════════════════


@app.get("/api/admin/analytics")
def admin_analytics(user: User = Depends(get_current_user)):
    """Platform-wide aggregate stats + growth data for the pitch deck."""
    if user.email != ADMIN_EMAIL:
        raise HTTPException(403, "Admin access only")

    from sqlalchemy import func
    from datetime import timedelta
    from database import CourseOutline

    db = SessionLocal()
    try:
        now = datetime.now(timezone.utc)
        thirty_days_ago = now - timedelta(days=30)
        cutoff_str = thirty_days_ago.strftime("%Y-%m-%d")

        total_users = db.query(User).count()
        total_messages = db.query(ChatMessage).count()
        total_notebooks = db.query(SavedNotebook).filter(SavedNotebook.is_premade == False).count()
        total_folders = db.query(StudyFolder).count()
        total_lessons_started = db.query(CourseOutline).count()
        total_sections_completed = (
            db.query(func.coalesce(func.sum(CourseOutline.current_section), 0)).scalar() or 0
        )

        avg_messages_per_user = round(total_messages / max(total_users, 1), 1)
        avg_notebooks_per_user = round(total_notebooks / max(total_users, 1), 1)

        users_with_lesson = db.query(CourseOutline.user_id).distinct().count()
        pct_started_lesson = round(users_with_lesson / max(total_users, 1) * 100, 1)

        users_completed = (
            db.query(CourseOutline.user_id)
            .filter(CourseOutline.current_section >= CourseOutline.total_sections)
            .filter(CourseOutline.total_sections > 0)
            .distinct()
            .count()
        )
        pct_completed_lesson = round(users_completed / max(total_users, 1) * 100, 1)

        lessons_with_sections = (
            db.query(CourseOutline).filter(CourseOutline.total_sections > 0).all()
        )
        if lessons_with_sections:
            avg_sections = round(
                sum(o.current_section for o in lessons_with_sections) / len(lessons_with_sections), 1
            )
        else:
            avg_sections = 0.0

        day_col = func.substr(User.created_at, 1, 10)
        signups_raw = (
            db.query(day_col.label("day"), func.count(User.id).label("cnt"))
            .filter(func.substr(User.created_at, 1, 10) >= cutoff_str)
            .group_by(day_col)
            .order_by(day_col)
            .all()
        )
        signups_per_day = [{"date": r.day, "count": r.cnt} for r in signups_raw]

        msg_day_col = func.substr(ChatMessage.created_at, 1, 10)
        messages_raw = (
            db.query(msg_day_col.label("day"), func.count(ChatMessage.id).label("cnt"))
            .filter(func.substr(ChatMessage.created_at, 1, 10) >= cutoff_str)
            .group_by(msg_day_col)
            .order_by(msg_day_col)
            .all()
        )
        messages_per_day = [{"date": r.day, "count": r.cnt} for r in messages_raw]

        total_feedback = db.query(UserFeedback).count()

        # ── Time on platform ──
        total_duration_ms = (
            db.query(func.coalesce(func.sum(ActivityEvent.duration_ms), 0)).scalar() or 0
        )
        total_hours = round(total_duration_ms / 3_600_000, 1)
        avg_hours_per_user = round(total_hours / max(total_users, 1), 2)

        feature_time_raw = (
            db.query(ActivityEvent.feature, func.sum(ActivityEvent.duration_ms).label("ms"))
            .group_by(ActivityEvent.feature)
            .all()
        )
        time_per_feature = {r.feature: round((r.ms or 0) / 3_600_000, 2) for r in feature_time_raw}

        # ── DAU / WAU / MAU ──
        today_str = now.strftime("%Y-%m-%d")
        seven_days_ago_str = (now - timedelta(days=7)).strftime("%Y-%m-%d")

        dau = (
            db.query(ActivityEvent.user_id)
            .filter(ActivityEvent.event_date == today_str)
            .distinct()
            .count()
        )
        wau = (
            db.query(ActivityEvent.user_id)
            .filter(ActivityEvent.event_date >= seven_days_ago_str)
            .distinct()
            .count()
        )
        mau = (
            db.query(ActivityEvent.user_id)
            .filter(ActivityEvent.event_date >= cutoff_str)
            .distinct()
            .count()
        )

        dau_trend_raw = (
            db.query(
                ActivityEvent.event_date.label("day"),
                func.count(func.distinct(ActivityEvent.user_id)).label("cnt"),
            )
            .filter(ActivityEvent.event_date >= cutoff_str)
            .group_by(ActivityEvent.event_date)
            .order_by(ActivityEvent.event_date)
            .all()
        )
        dau_trend = [{"date": r.day, "count": r.cnt} for r in dau_trend_raw]

        # ── Retention ──
        all_users = db.query(User.id, func.substr(User.created_at, 1, 10).label("signup")).all()
        user_activity_dates = {}
        activity_rows = (
            db.query(ActivityEvent.user_id, ActivityEvent.event_date)
            .distinct()
            .all()
        )
        for uid, ed in activity_rows:
            user_activity_dates.setdefault(uid, set()).add(ed)

        retention = {"day_1": 0.0, "day_7": 0.0, "day_30": 0.0}
        if all_users:
            for offset_days, key in [(1, "day_1"), (7, "day_7"), (30, "day_30")]:
                eligible = 0
                returned = 0
                for uid, signup in all_users:
                    if not signup:
                        continue
                    try:
                        signup_dt = datetime.strptime(signup, "%Y-%m-%d")
                    except (ValueError, TypeError):
                        continue
                    target = (signup_dt + timedelta(days=offset_days)).strftime("%Y-%m-%d")
                    if target > today_str:
                        continue
                    eligible += 1
                    if uid in user_activity_dates and target in user_activity_dates[uid]:
                        returned += 1
                retention[key] = round(returned / max(eligible, 1) * 100, 1)

        return {
            "headline": {
                "total_users": total_users,
                "total_messages": total_messages,
                "total_notebooks": total_notebooks,
                "total_folders": total_folders,
                "total_lessons_started": total_lessons_started,
                "total_sections_completed": int(total_sections_completed),
                "total_feedback": total_feedback,
            },
            "engagement": {
                "avg_messages_per_user": avg_messages_per_user,
                "avg_notebooks_per_user": avg_notebooks_per_user,
                "pct_started_lesson": pct_started_lesson,
                "pct_completed_lesson": pct_completed_lesson,
                "avg_sections_per_lesson": avg_sections,
            },
            "growth": {
                "signups_per_day": signups_per_day,
                "messages_per_day": messages_per_day,
            },
            "time": {
                "total_hours": total_hours,
                "avg_hours_per_user": avg_hours_per_user,
                "per_feature": time_per_feature,
            },
            "active_users": {
                "dau": dau,
                "wau": wau,
                "mau": mau,
                "dau_trend": dau_trend,
            },
            "retention": retention,
        }
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# FEEDBACK
# ═══════════════════════════════════════════════════════════════════════════


class FeedbackRequest(BaseModel):
    category: str = "other"
    message: str
    page: str = ""


@app.post("/api/feedback")
def submit_feedback(body: FeedbackRequest, user: User = Depends(get_current_user)):
    """Any authenticated user can submit feedback."""
    if not body.message.strip():
        raise HTTPException(400, "Message cannot be empty")

    db = SessionLocal()
    try:
        fb = UserFeedback(
            user_id=user.id,
            category=body.category,
            message=body.message.strip(),
            page=body.page,
        )
        db.add(fb)
        db.commit()
        return {"ok": True, "id": fb.id}
    finally:
        db.close()


@app.get("/api/admin/feedback")
def admin_feedback(user: User = Depends(get_current_user)):
    """Admin-only: return all user feedback with user info."""
    if user.email != ADMIN_EMAIL:
        raise HTTPException(403, "Admin access only")

    db = SessionLocal()
    try:
        rows = (
            db.query(UserFeedback, User)
            .join(User, UserFeedback.user_id == User.id)
            .order_by(UserFeedback.created_at.desc())
            .all()
        )
        return {
            "feedback": [
                {
                    "id": fb.id,
                    "user_name": u.name,
                    "user_email": u.email,
                    "category": fb.category,
                    "message": fb.message,
                    "page": fb.page,
                    "created_at": fb.created_at.isoformat() if fb.created_at else None,
                }
                for fb, u in rows
            ]
        }
    finally:
        db.close()


# ═══════════════════════════════════════════════════════════════════════════
# LIVE PRESENCE
# ═══════════════════════════════════════════════════════════════════════════


@app.post("/api/heartbeat")
def heartbeat(user: User = Depends(get_current_user)):
    """Called every ~30s by each active browser tab."""
    _live_users[user.id] = {
        "name": user.name,
        "email": user.email,
        "last_seen": datetime.now(timezone.utc),
    }
    return {"ok": True}


@app.get("/api/admin/live-users")
def admin_live_users(user: User = Depends(get_current_user)):
    """Return currently active users (heartbeat within last 60s). Admin only."""
    if user.email != ADMIN_EMAIL:
        raise HTTPException(403, "Admin access only")

    now = datetime.now(timezone.utc)
    active = []
    stale_ids = []
    for uid, info in _live_users.items():
        age = (now - info["last_seen"]).total_seconds()
        if age <= HEARTBEAT_TIMEOUT:
            active.append({
                "user_id": uid,
                "name": info["name"],
                "email": info["email"],
                "seconds_ago": int(age),
            })
        else:
            stale_ids.append(uid)

    for uid in stale_ids:
        del _live_users[uid]

    return {"count": len(active), "users": active}


# ═══════════════════════════════════════════════════════════════════════════
# HEALTH
# ═══════════════════════════════════════════════════════════════════════════

@app.get("/api/health")
def health():
    db = SessionLocal()
    try:
        paper_count = db.query(Paper).count()
        user_count = db.query(User).count()
        return {"status": "ok", "papers": paper_count, "users": user_count}
    finally:
        db.close()


def _get_paper_paths(course: str | None = None) -> list[Path]:
    """Get past paper JSON paths, optionally filtered by course.

    If course is None, returns nothing — we only match papers
    when we know the course to avoid cross-course contamination.
    """
    if not course:
        return []
    paper_files = []
    if PAPERS_DIR.exists():
        for f in PAPERS_DIR.glob("*.json"):
            if f.name in ("notebooks.json", "notebookContent.json"):
                continue
            try:
                import json as _json
                with open(f) as _fh:
                    data = _json.load(_fh)
                if data.get("course", "").upper() == course.upper():
                    paper_files.append(f)
            except Exception:
                pass
    return paper_files


if __name__ == "__main__":
    import uvicorn

    print("=" * 50)
    print("  Coast API Server")
    print("=" * 50)
    print(f"  Papers dir:  {PAPERS_DIR}")
    print(f"  Generated:   {GENERATED_DIR}")
    print(f"  Database:    {Path(__file__).parent / 'coast.db'}")
    print("=" * 50)
    uvicorn.run(app, host="0.0.0.0", port=8000)
